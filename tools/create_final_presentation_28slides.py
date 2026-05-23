from __future__ import annotations

import math
import re
import shutil
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs"
ASSETS = OUT / "ppt_thesis_visual_assets" / "figures"
PPTX_PATH = OUT / "final_presentation_maintiq_predict_28slides.pptx"
PREVIEW_DIR = OUT / "final_presentation_maintiq_predict_28slides_preview"
CONTACT_SHEET = PREVIEW_DIR / "contact_sheet.png"
NOTES_PATH = OUT / "final_presentation_28slides_speaker_notes.md"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Malgun Gothic"


class C:
    GREEN = "0F6B4F"
    GREEN_DARK = "064B36"
    GREEN_PALE = "E7F3EF"
    BLUE = "2C6EA3"
    BLUE_PALE = "EAF3FA"
    ORANGE = "D97918"
    ORANGE_PALE = "FFF2E2"
    RED = "C0392B"
    INK = "111827"
    MUTED = "5B6472"
    LINE = "CFDCE6"
    BG = "FFFFFF"
    SOFT_BG = "F6F9FB"
    WHITE = "FFFFFF"


def rgb(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def pil_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


@dataclass
class SlideSpec:
    title: str
    message: str
    bullets: list[str]
    asset: str | None
    asset_caption: str
    talk: str
    qa: str
    layout: str = "standard"


SLIDES: list[SlideSpec] = [
    SlideSpec(
        "MaintiQ Predict",
        "ML 예측, Predictive SPC, GenAI 리포트, 승인형 작업지시를 결합한 스마트 제조 예지보전 시스템",
        ["전남대학교 산업공학과", "최종발표", "CSV 기반 예측에서 작업지시 의사결정까지 연결"],
        None,
        "",
        "안녕하십니까. 본 발표에서는 ML 예측, Predictive SPC, GenAI 리포트, 승인형 작업지시를 결합한 스마트 제조 예지보전 시스템 MaintiQ Predict를 설명드리겠습니다.",
        "제품명은 구현체이고, 연구 핵심은 예측 결과를 운영 의사결정으로 연결한 통합 구조입니다.",
        "cover",
    ),
    SlideSpec(
        "목차",
        "발표 흐름은 문제 정의에서 시스템 구현, 실험 결과, 결론으로 이어진다.",
        ["Motivation", "Contributions", "참고문헌 Review", "Main Part", "실험 설계", "실험 결과", "결론 및 References"],
        None,
        "",
        "발표는 제조 설비 고장 문제와 기존 보전 방식의 한계를 설명하고, 제안 시스템과 실험 결과, 주장 경계를 말씀드리는 순서로 진행하겠습니다.",
        "예시 최종발표 흐름에 맞춰 7개 큰 흐름으로 구성했습니다.",
        "toc",
    ),
    SlideSpec(
        "Motivation 1: 제조 설비 고장 문제",
        "설비 고장은 downtime, 품질 손실, 정비 비용 증가로 이어진다.",
        ["고장은 생산 중단과 납기 지연으로 연결된다.", "센서 데이터는 많지만 의사결정 흐름과 분리되기 쉽다.", "고장 가능성을 사전에 판단하고 조치 검토로 연결해야 한다."],
        None,
        "",
        "제조 설비 고장은 단순한 장비 문제가 아니라 생산 중단, 품질 손실, 납기 지연으로 이어집니다. 따라서 고장 가능성을 사전에 파악하고 의사결정으로 연결하는 구조가 필요합니다.",
        "기본 모델은 AI4I, 공개 산업 검증은 SCANIA 등 공개 benchmark를 사용했습니다.",
        "bullets",
    ),
    SlideSpec(
        "Motivation 2: 기존 보전 방식 한계",
        "사후보전과 예방보전은 놓치는 고장과 과잉 정비라는 trade-off를 갖는다.",
        ["사후보전은 고장 후 대응이라 downtime이 커질 수 있다.", "예방보전은 주기 기반이라 과잉 정비 가능성이 있다.", "예지보전은 데이터 품질과 검증이 전제되어야 한다."],
        "04_maintenance_strategy_comparison_table.png",
        "보전 전략 비교표",
        "사후보전은 고장 후 대응이므로 downtime이 커질 수 있고, 예방보전은 정해진 주기에 따라 정비해 과잉 정비가 발생할 수 있습니다. 본 연구는 센서 데이터 기반 예지보전을 통해 이 간격을 줄이는 방향을 다룹니다.",
        "예지보전도 false alarm과 missed failure trade-off를 함께 봐야 합니다.",
    ),
    SlideSpec(
        "Contributions",
        "본 연구의 기여는 단일 모델 성능보다 운영 흐름 통합에 있다.",
        ["XGBoost 기반 고장 확률 예측과 threshold tuning", "Predictive SPC로 위험 확률 흐름 모니터링", "GenAI 관리자 참고 리포트 생성", "승인형 작업지시 workflow와 공개 benchmark 검증"],
        None,
        "",
        "본 연구의 차별점은 예측 모델 하나에 그치지 않고, threshold 정책, Predictive SPC, GenAI 리포트, 승인형 작업지시를 하나의 제품형 MVP로 구현했다는 점입니다.",
        "Something new는 모델 자체보다 예측 결과를 운영 의사결정 흐름으로 연결한 통합성입니다.",
        "cards",
    ),
    SlideSpec(
        "Review 1: 예지보전 / CBM",
        "CBM 선행연구는 상태 기반 정비의 필요성을 제시한다.",
        ["Jardine et al. (2006): 기계 진단과 예후 기반 CBM 정리", "Carvalho et al. (2019): ML 기반 예지보전 연구 동향", "Montgomery (2019): SPC 관리도와 품질관리 이론"],
        None,
        "",
        "선행연구는 상태 기반 정비와 예지보전의 필요성을 강조합니다. 본 연구는 이 흐름을 따라가되, 예측 결과를 관리자 리포트와 작업지시로 연결하는 구현에 초점을 두었습니다.",
        "산업공학 지식은 SPC, 보전 전략 비교, FMEA/RPN식 위험 우선순위, cost simulation에 반영했습니다.",
        "review",
    ),
    SlideSpec(
        "Review 2: ML 기반 예지보전",
        "불균형 데이터에서는 PR-AUC와 threshold 정책이 중요하다.",
        ["AI4I 2020 / Matzka: 예지보전 실험 데이터", "Chen & Guestrin (2016): XGBoost", "Chawla et al. (2002): SMOTE"],
        None,
        "",
        "ML 기반 예지보전에서는 고장 클래스가 매우 적은 불균형 문제가 중요합니다. 따라서 본 연구는 PR-AUC, recall, precision, threshold tuning을 함께 확인했습니다.",
        "고장이 적은 데이터에서는 accuracy가 높아도 고장 탐지가 약할 수 있어 PR-AUC와 recall이 중요합니다.",
        "review",
    ),
    SlideSpec(
        "Review 3: SPC / SHAP / Cost-sensitive Learning",
        "예측 결과는 관리도, 설명성, 비용 민감 의사결정과 결합될 때 운영 가치가 커진다.",
        ["SPC: 위험 확률 흐름과 관리한계", "SHAP: 주요 위험 요인 설명", "Elkan (2001), SCANIA: 비용 민감 평가와 official cost metric"],
        None,
        "",
        "SPC는 위험 추세와 관리한계를 보는 관점이고, SHAP은 위험 요인을 설명하며, cost-sensitive learning은 false alarm과 missed failure의 비용 차이를 반영합니다.",
        "본 연구의 비용 수치는 실제 원화 절감이 아니라 공식 cost metric과 normalized simulation입니다.",
        "review",
    ),
    SlideSpec(
        "Main Part: 전체 시스템 구조",
        "CSV 입력부터 예측, SPC, 리포트, 작업지시까지 연결된다.",
        ["입력 데이터와 예측 결과가 운영 화면에서 연결된다.", "GenAI 리포트는 관리자 참고용이다.", "작업지시는 사람이 승인/검토/반려한다."],
        "01_system_architecture.png",
        "전체 시스템 아키텍처",
        "제안 시스템은 센서 CSV를 입력받아 전처리와 예측을 수행하고, 위험 확률을 SPC 관점으로 모니터링하며, GenAI 리포트와 승인형 작업지시로 이어집니다.",
        "현재는 CSV와 로컬 기반 MVP이며 실제 설비망 연결은 향후 과제입니다.",
        "image_full",
    ),
    SlideSpec(
        "Main Part: MaintiQ Predict 제품 화면",
        "운영자는 데스크톱 앱에서 CSV 예측, 위험 확인, 리포트, 작업지시를 처리한다.",
        ["사용자 앱은 연구 용어를 숨긴 제품형 화면이다.", "Admin 콘솔은 연구/검증 자료 확인용이다."],
        "20_app_main_screen.png",
        "MaintiQ Predict 메인 화면",
        "이 화면은 사용자용 데스크톱 앱입니다. 연구 용어를 숨기고 운영자가 바로 사용할 수 있도록 데이터 예측, 위험 모니터링, AI 리포트, 작업지시 메뉴로 구성했습니다.",
        "사용자 앱은 PySide6 데스크톱이고, 연구 검증은 별도 Admin 콘솔에서 봅니다.",
    ),
    SlideSpec(
        "Main Part: 데이터 입력 및 전처리",
        "AI4I 데이터는 ID와 leakage 컬럼을 제거하고 Type을 one-hot encoding한다.",
        ["제거: UDI, Product ID", "누수 방지: TWF/HDF/PWF/OSF/RNF 제거", "Type one-hot encoding, stratified split 적용"],
        "02_data_preprocessing_pipeline.png",
        "데이터 입력 및 전처리 흐름",
        "학습에서는 UDI와 Product ID 같은 식별 컬럼을 제거했고, TWF/HDF/PWF/OSF/RNF처럼 target과 직접 연결되는 leakage 컬럼도 제거했습니다. Type은 one-hot encoding했습니다.",
        "failure-type 컬럼은 고장 원인 라벨에 가까워 실제 예측 시 데이터 누수를 만들 수 있기 때문에 제거했습니다.",
    ),
    SlideSpec(
        "Main Part: 예측 모델 구조",
        "Logistic Regression을 기준선으로 두고 XGBoost를 주 예측 모델로 비교했다.",
        ["Baseline: Logistic Regression", "Main model: XGBoost", "평가 지표: precision, recall, F1, ROC-AUC, PR-AUC"],
        "26_shap_summary.png",
        "SHAP summary: 주요 위험 요인",
        "baseline으로 Logistic Regression을 사용했고, 비선형 센서 조합을 반영하기 위해 XGBoost를 주 모델로 사용했습니다. 성능 비교는 precision, recall, F1, ROC-AUC, PR-AUC로 수행했습니다.",
        "AI4I 기준 주 모델은 XGBoost이며 PR-AUC가 0.8014입니다.",
    ),
    SlideSpec(
        "Main Part: Threshold Policy",
        "고정 0.5 기준 대신 F1 기준 threshold 0.87을 선택했다.",
        ["0.05~0.95 범위에서 threshold search", "선택 기준: F1-score", "선택 결과: threshold 0.87, F1 0.7752"],
        "06_threshold_tuning_table.png",
        "Threshold tuning 결과표",
        "예측 확률은 0.5 기준으로만 판단하지 않았습니다. threshold를 0.05부터 0.95까지 탐색했고, F1 기준으로 0.87을 선택했습니다.",
        "0.87은 해당 test split에서 precision과 recall의 균형인 F1-score가 가장 높았기 때문입니다.",
    ),
    SlideSpec(
        "Main Part: Predictive SPC",
        "예측 확률을 시계열 관리도처럼 보아 위험 흐름을 관리한다.",
        ["ML 위험 확률을 시간 흐름으로 배열", "관리한계와 high-risk point 표시", "SPC-only rule과 ML+SPC 정책 비교"],
        "24_spc_risk_chart.png",
        "Predictive SPC 위험 확률 chart",
        "Predictive SPC는 단일 행 예측을 넘어서 위험 확률의 흐름을 보는 구조입니다. 고위험 구간과 관리한계를 함께 보여주어 운영자가 위험 추세를 판단할 수 있게 했습니다.",
        "기존 SPC와 달리 단일 센서 control limit만 보는 것이 아니라 ML 위험 확률을 함께 사용합니다.",
    ),
    SlideSpec(
        "Main Part: GenAI 관리자 리포트",
        "Gemini API가 위험 context를 관리자 참고 리포트로 요약한다.",
        ["mode: gemini_generate_content:gemini-2.5-flash", "probability: 0.993616, threshold: 0.87", "상태: High Risk, 주요 요인: torque, speed, air temperature"],
        "08_genai_report_evidence_table.png",
        "Gemini AI 리포트 검증표",
        "Gemini generateContent API를 사용해 예측 확률 0.993616, 기준 0.87, High Risk 상태를 관리자 참고 리포트로 변환했습니다. 이 리포트는 조치 검토를 돕는 설명 자료입니다.",
        "AI 리포트는 참고용이고 최종 결정은 승인형 작업지시에서 사람이 기록합니다.",
    ),
    SlideSpec(
        "Main Part: 승인형 작업지시 Workflow",
        "예측 결과는 자동 실행이 아니라 작업자 승인/검토/반려로 이어진다.",
        ["센서 이벤트 생성", "작업지시 초안 생성", "작업자 결정과 이력 저장"],
        "03_work_order_workflow.png",
        "승인형 작업지시 workflow",
        "시스템은 고위험 이벤트를 만들고 작업지시 초안을 생성하지만, 정비를 자동 실행하지 않습니다. 작업자가 승인, 검토 필요, 반려를 선택하고 이력이 저장됩니다.",
        "자동 제어가 아니라 사람 승인 구조라 안전한 MVP 경계에 맞습니다.",
    ),
    SlideSpec(
        "실험 설계",
        "AI4I, 공개 benchmark, cost simulation을 분리해 검증했다.",
        ["AI4I baseline과 threshold tuning", "SMOTE/전략 비교와 SPC 정책 비교", "SCANIA official cost metric과 public benchmark 확장"],
        "11_experiment_design_table.png",
        "실험 설계 요약표",
        "실험은 AI4I baseline, threshold tuning, SMOTE 비교, SPC-only vs ML+SPC, cost simulation, SCANIA official cost metric으로 나누어 진행했습니다.",
        "실제 회사 실증은 별도 로그가 필요하며, 본 연구는 공개 데이터와 simulation 검증입니다.",
    ),
    SlideSpec(
        "실험 결과 1: AI4I Baseline",
        "XGBoost가 PR-AUC 0.8014로 Logistic Regression보다 우수했다.",
        ["XGBoost PR-AUC: 0.8014", "XGBoost ROC-AUC: 0.9736", "불균형 데이터이므로 PR-AUC를 강조"],
        "05_ai4i_baseline_table.png",
        "AI4I baseline 성능표",
        "AI4I 기준 XGBoost는 PR-AUC 0.8014, ROC-AUC 0.9736을 기록했습니다. Logistic Regression보다 고장 탐지 ranking 성능이 높았습니다.",
        "고장이 적은 불균형 데이터에서는 positive class 탐지 성능을 더 직접적으로 보여주는 PR-AUC가 중요합니다.",
    ),
    SlideSpec(
        "실험 결과 2: Threshold Tuning",
        "Threshold 0.87에서 F1-score 0.7752를 얻었다.",
        ["Default 0.5 F1: 0.5911", "Selected 0.87 F1: 0.7752", "precision 0.8197, recall 0.7353"],
        "23_threshold_tuning.png",
        "Threshold tuning plot",
        "기본 threshold 0.5와 비교해 F1 기준 threshold 0.87을 선택했습니다. 이때 precision은 0.8197, recall은 0.7353, F1은 0.7752입니다.",
        "false alarm을 줄이는 대신 missed failure가 늘 수 있어 운영 정책에 따라 threshold를 조정해야 합니다.",
    ),
    SlideSpec(
        "실험 결과 3: SMOTE / 모델 전략 비교",
        "SMOTE는 무조건 개선이 아니라 precision-recall trade-off로 해석해야 한다.",
        ["Logistic Regression, XGBoost, SMOTE 조합 비교", "Threshold tuning 조합 비교", "PR curve 기반 trade-off 확인"],
        "33_model_strategy_pr_curve.png",
        "모델 전략 PR curve",
        "불균형 처리를 위해 SMOTE 적용 모델도 비교했습니다. 결과는 SMOTE가 항상 우수하다는 식이 아니라, precision과 recall의 trade-off로 해석했습니다.",
        "SMOTE는 비교 실험 근거로 사용하고, 최종 운영 판단은 지표별 trade-off를 보고 결정합니다.",
    ),
    SlideSpec(
        "실험 결과 4: SPC-only vs ML+SPC",
        "SPC-only는 alert가 적고 recall이 낮았으며, ML+SPC는 recall을 높였다.",
        ["SPC-only F1: 0.1600", "ML threshold F1: 0.7752", "ML+SPC F1: 0.7051, recall: 0.8088"],
        "07_spc_vs_ml_table.png",
        "SPC-only vs ML+SPC 비교표",
        "SPC-only torque rule은 alert가 7개로 적지만 recall이 0.0882였습니다. ML threshold는 F1 0.7752, ML+SPC는 recall 0.8088로 더 적극적인 탐지 특성을 보였습니다.",
        "F1은 ML threshold가 높고, recall-first 운영이면 ML+SPC를 고려할 수 있습니다.",
    ),
    SlideSpec(
        "실험 결과 5: Operational Cost Simulation",
        "false alarm과 missed failure 비용 가중치로 운영 정책을 비교했다.",
        ["실제 원화 비용이 아니라 normalized cost", "conservative, balanced, high_downtime 시나리오", "정책별 false alarm / missed failure trade-off"],
        "28_operational_value_simulation.png",
        "Operational value simulation",
        "운영 비용은 실제 원화가 아니라 normalized cost simulation으로 평가했습니다. false alarm과 missed failure의 상대 비용을 두고 정책별 의사결정 가능성을 비교했습니다.",
        "이는 실제 비용 실증이 아니라 비용 민감 정책 비교 시뮬레이션입니다.",
    ),
    SlideSpec(
        "실험 결과 6: SCANIA Official Cost Metric",
        "SCANIA 공개 benchmark에서 official cost metric 기준 rule 대비 17.02% 개선 가능성을 확인했다.",
        ["SCANIA Component X 공개 산업 데이터", "official class 0~4 cost matrix", "XGBoost cost-optimized: rule 대비 17.02% 개선"],
        "09_scania_cost_metric_table.png",
        "SCANIA official cost metric 비교표",
        "SCANIA Component X는 실제 SCANIA fleet에서 수집된 공개 benchmark입니다. official cost matrix 기준 XGBoost cost-optimized 전략이 rule baseline 대비 17.02% 낮은 cost metric을 보였습니다.",
        "이는 실제 회사 비용 절감이 아니라 SCANIA official cost metric 기준 공개 benchmark 결과입니다.",
    ),
    SlideSpec(
        "실험 결과 7: Public Benchmark Extension",
        "MetroPT-3, C-MAPSS, IMS/FEMTO는 확장 검증 경로로 정리했다.",
        ["MetroPT-3: compressor anomaly/failure horizon", "NASA C-MAPSS: RUL / lead-time", "IMS/FEMTO: bearing run-to-failure"],
        "31_public_benchmark_cost_chart.png",
        "Public benchmark simulated cost summary",
        "AI4I와 SCANIA 외에도 MetroPT-3, NASA C-MAPSS, IMS/FEMTO 데이터를 고려할 수 있는 adapter와 검증 산출물을 정리했습니다. 이는 실제 회사 실증 전 공개 benchmark 확장 근거입니다.",
        "원본 데이터 유무에 따라 sample smoke와 full benchmark가 구분됩니다.",
    ),
    SlideSpec(
        "제품 구현 결과: Full/Lite 데스크톱 앱",
        "정밀 분석 모드와 빠른 점검 모드를 분리해 설치성과 분석성을 모두 고려했다.",
        ["Full: 정밀 분석 모드, XGBoost/SHAP 기반", "Lite: 빠른 점검 모드, 경량 운영 점수", "두 모드는 목적이 달라 결과가 다를 수 있음"],
        "21_app_lite_screen.png",
        "Lite 데스크톱 앱 화면",
        "Full은 XGBoost/SHAP 기반 정밀 분석 모드이고, Lite는 작은 설치본과 경량 운영 점수를 사용하는 빠른 점검 모드입니다. 두 모드는 목적이 달라 결과가 다를 수 있습니다.",
        "연구 검증과 일반 배포의 의존성/용량 요구가 다르기 때문에 두 버전으로 나눴습니다.",
    ),
    SlideSpec(
        "Claim Boundary: 가능한 주장 / 피해야 할 표현",
        "본 연구는 공개 데이터와 로컬 MVP 검증이며 실제 현장 실증은 별도 데이터가 필요하다.",
        ["가능: 공개 데이터 기반 모델/정책 비교", "가능: SCANIA official cost metric 개선", "피해야 함: 실제 현장 비용 절감 또는 PLC 운영망 배포처럼 표현"],
        "10_claim_boundary_table.png",
        "가능한 주장 / 피해야 할 표현",
        "발표에서 가장 중요한 경계는 과장하지 않는 것입니다. 공개 benchmark와 로컬 simulation은 주장할 수 있지만, 실제 회사 비용 절감이나 PLC 운영망 배포는 주장하지 않습니다.",
        "실제 현장 실증에는 labeled sensor CSV, maintenance history, downtime/cost log가 필요합니다.",
    ),
    SlideSpec(
        "결론 및 향후 연구",
        "예측에서 운영 의사결정까지 연결한 제품형 MVP를 구현했다.",
        ["구현: 예측, SPC, GenAI, 작업지시 통합", "검증: AI4I, SCANIA, public benchmark", "향후: 실제 회사 로그 기반 field validation"],
        "12_paper_structure_29p.png",
        "논문 29쪽 내외 구성",
        "결론적으로 본 연구는 고장 확률 예측, SPC 모니터링, GenAI 리포트, 승인형 작업지시를 연결한 제품형 MVP를 구현했습니다. 향후 실제 회사 로그로 성능과 비용 효과를 검증하는 것이 다음 단계입니다.",
        "CSV 기반 MVP로는 시연 가능하지만, 현장 운영 적용은 데이터 계약과 설비 연동 검증이 필요합니다.",
    ),
    SlideSpec(
        "References",
        "핵심 참고문헌은 예지보전, SPC, ML, 설명가능성, 공개 benchmark를 포괄한다.",
        ["Jardine et al. (2006), Carvalho et al. (2019), Montgomery (2019)", "AI4I 2020, SCANIA Component X, MetroPT-3, NASA C-MAPSS, IMS", "XGBoost, SMOTE, SHAP, cost-sensitive learning, Gemini API"],
        None,
        "",
        "참고문헌은 예지보전/CBM, SPC, XGBoost, SMOTE, SHAP, SCANIA Component X, AI4I 2020, Gemini API 문서를 중심으로 구성했습니다.",
        "이론, 데이터셋, 모델, 구현 API를 모두 포함하도록 구성했습니다.",
        "references",
    ),
]


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 16, color: str = C.INK, bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(C.INK)
        p.space_after = Pt(10)


def add_shape(slide, x: float, y: float, w: float, h: float, fill: str = C.WHITE, line: str = C.LINE, rounded: bool = True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(0.8)
    return shp


def add_footer(slide, page_no: int) -> None:
    add_text(slide, "전남대학교 산업공학과 | MaintiQ Predict", 0.35, 7.15, 5.2, 0.18, 8, C.MUTED)
    add_text(slide, str(page_no), 12.55, 7.15, 0.35, 0.18, 8, C.MUTED, align=PP_ALIGN.RIGHT)


def add_title(slide, title: str, page_no: int) -> None:
    add_text(slide, title, 0.62, 0.38, 11.6, 0.46, 25, C.GREEN_DARK, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.05), Inches(11.95), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(C.GREEN)
    line.line.color.rgb = rgb(C.GREEN)
    add_footer(slide, page_no)


def add_image_fit(slide, image_name: str, x: float, y: float, w: float, h: float) -> None:
    path = ASSETS / image_name
    if not path.exists():
        add_shape(slide, x, y, w, h, C.SOFT_BG, C.LINE)
        add_text(slide, f"이미지 없음: {image_name}", x + 0.15, y + 0.15, w - 0.3, 0.4, 13, C.RED, True)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / (iw / 96), h / (ih / 96))
    pic_w = iw / 96 * scale
    pic_h = ih / 96 * scale
    px = x + (w - pic_w) / 2
    py = y + (h - pic_h) / 2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(pic_w), height=Inches(pic_h))


def build_cover(slide, spec: SlideSpec) -> None:
    add_text(slide, "MaintiQ Predict", 0.72, 1.12, 8.5, 0.7, 36, C.GREEN_DARK, True)
    add_text(slide, spec.message, 0.76, 2.0, 10.7, 0.75, 20, C.INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.76), Inches(2.9), Inches(11.7), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(C.GREEN)
    line.line.color.rgb = rgb(C.GREEN)
    for idx, txt in enumerate(spec.bullets):
        add_text(slide, txt, 0.9, 3.45 + idx * 0.48, 8.5, 0.35, 16, C.MUTED if idx else C.INK, idx == 0)
    add_shape(slide, 9.6, 3.25, 2.65, 1.75, C.GREEN_PALE, C.GREEN)
    add_text(slide, "Prediction\nSPC\nGenAI\nWork-order", 9.88, 3.55, 2.2, 1.0, 20, C.GREEN_DARK, True, PP_ALIGN.CENTER)
    add_footer(slide, 1)


def build_toc(slide, spec: SlideSpec, page_no: int) -> None:
    add_title(slide, "목차", page_no)
    labels = spec.bullets
    for i, label in enumerate(labels):
        y = 1.65 + i * 0.65
        add_text(slide, f"{i + 1:02d}", 1.0, y, 0.65, 0.34, 18, C.GREEN, True)
        add_text(slide, label, 1.85, y, 8.8, 0.34, 19, C.INK, True)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.85), Inches(y + 0.42), Inches(9.4), Inches(0.008))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(C.LINE)
        line.line.color.rgb = rgb(C.LINE)


def build_cards(slide, spec: SlideSpec, page_no: int) -> None:
    add_title(slide, spec.title, page_no)
    add_text(slide, spec.message, 0.8, 1.35, 11.8, 0.38, 18, C.MUTED)
    x_positions = [0.8, 3.95, 7.1, 10.25]
    colors = [C.GREEN_PALE, C.BLUE_PALE, C.ORANGE_PALE, C.GREEN_PALE]
    for i, bullet in enumerate(spec.bullets[:4]):
        add_shape(slide, x_positions[i], 2.25, 2.5, 2.75, colors[i], C.LINE)
        add_text(slide, f"{i + 1}", x_positions[i] + 0.18, 2.45, 0.45, 0.4, 20, C.GREEN_DARK, True)
        add_text(slide, bullet, x_positions[i] + 0.25, 3.05, 2.0, 1.3, 17, C.INK, True)


def build_standard(slide, spec: SlideSpec, page_no: int) -> None:
    add_title(slide, spec.title, page_no)
    add_text(slide, spec.message, 0.72, 1.32, 11.9, 0.42, 17, C.MUTED)
    if spec.asset:
        add_bullets(slide, spec.bullets, 0.78, 2.05, 4.3, 3.8, 16)
        add_shape(slide, 5.35, 1.83, 7.15, 4.9, C.WHITE, C.LINE)
        add_image_fit(slide, spec.asset, 5.55, 2.03, 6.75, 4.5)
        add_text(slide, spec.asset_caption, 5.55, 6.5, 6.8, 0.22, 9, C.MUTED, align=PP_ALIGN.CENTER)
    else:
        add_bullets(slide, spec.bullets, 0.9, 2.0, 10.8, 4.2, 18)


def build_image_full(slide, spec: SlideSpec, page_no: int) -> None:
    add_title(slide, spec.title, page_no)
    add_text(slide, spec.message, 0.72, 1.28, 11.9, 0.35, 16, C.MUTED)
    if spec.asset:
        add_image_fit(slide, spec.asset, 0.85, 1.72, 11.8, 5.05)


def build_review(slide, spec: SlideSpec, page_no: int) -> None:
    add_title(slide, spec.title, page_no)
    add_text(slide, spec.message, 0.75, 1.3, 11.8, 0.38, 17, C.MUTED)
    for i, bullet in enumerate(spec.bullets):
        y = 2.05 + i * 1.15
        add_text(slide, f"Review {i + 1}", 0.95, y, 1.1, 0.25, 10, C.GREEN, True)
        add_shape(slide, 2.05, y - 0.08, 9.9, 0.72, C.SOFT_BG, C.LINE)
        add_text(slide, bullet, 2.25, y + 0.08, 9.45, 0.38, 16, C.INK)


def build_references(slide, spec: SlideSpec, page_no: int) -> None:
    add_title(slide, spec.title, page_no)
    refs = [
        "Jardine, Lin, & Banjevic (2006). Machinery diagnostics and prognostics implementing CBM.",
        "Carvalho et al. (2019). ML methods applied to predictive maintenance: systematic review.",
        "Montgomery (2019). Introduction to Statistical Quality Control.",
        "Chen & Guestrin (2016). XGBoost: A scalable tree boosting system.",
        "Chawla et al. (2002). SMOTE.",
        "Lundberg & Lee (2017). SHAP.",
        "AI4I 2020 Dataset; SCANIA Component X Dataset; MetroPT-3; NASA C-MAPSS; IMS Bearings.",
        "Google Gemini generateContent API; OpenAI Responses API; PyInstaller.",
    ]
    add_bullets(slide, refs, 0.75, 1.35, 11.8, 5.5, 13)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for idx, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(C.BG)
        if spec.layout == "cover":
            build_cover(slide, spec)
        elif spec.layout == "toc":
            build_toc(slide, spec, idx)
        elif spec.layout == "cards":
            build_cards(slide, spec, idx)
        elif spec.layout == "image_full":
            build_image_full(slide, spec, idx)
        elif spec.layout == "review":
            build_review(slide, spec, idx)
        elif spec.layout == "references":
            build_references(slide, spec, idx)
        else:
            build_standard(slide, spec, idx)
    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.ImageFont, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    line = ""
    for word in words:
        cand = word if not line else f"{line} {word}"
        if draw.textbbox((0, 0), cand, font=fnt)[2] <= max_width:
            line = cand
        else:
            if line:
                lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines or [""]


def pil_color(hex_value: str) -> str:
    return "#" + hex_value.lstrip("#")


def preview_slide(spec: SlideSpec, idx: int) -> Image.Image:
    img = Image.new("RGB", (1600, 900), pil_color(C.BG))
    draw = ImageDraw.Draw(img)
    title_font = pil_font(36, True)
    body_font = pil_font(23)
    small_font = pil_font(17)
    if spec.layout == "cover":
        draw.text((88, 130), "MaintiQ Predict", fill="#" + C.GREEN_DARK, font=pil_font(54, True))
        draw.text((92, 245), "ML 예측, Predictive SPC, GenAI 리포트, 승인형 작업지시를 결합한 스마트 제조 예지보전 시스템", fill="#" + C.INK, font=pil_font(26))
        draw.rectangle((92, 342, 1470, 348), fill="#" + C.GREEN)
        draw.text((108, 420), "전남대학교 산업공학과 · 최종발표", fill="#" + C.MUTED, font=pil_font(23, True))
    else:
        draw.text((75, 46), spec.title, fill="#" + C.GREEN_DARK, font=title_font)
        draw.rectangle((75, 124, 1515, 129), fill="#" + C.GREEN)
        y = 160
        for line in wrap_text(draw, spec.message, body_font, 1430)[:2]:
            draw.text((78, y), line, fill="#" + C.MUTED, font=body_font)
            y += 34
        if spec.asset:
            asset_path = ASSETS / spec.asset
            if asset_path.exists():
                pic = Image.open(asset_path).convert("RGB")
                pic.thumbnail((820, 520))
                img.paste(pic, (720, 250))
            yb = 250
            for bullet in spec.bullets[:5]:
                for line in wrap_text(draw, "• " + bullet, body_font, 590)[:2]:
                    draw.text((95, yb), line, fill="#" + C.INK, font=body_font)
                    yb += 34
                yb += 18
        else:
            yb = 240
            for bullet in spec.bullets[:8]:
                for line in wrap_text(draw, "• " + bullet, body_font, 1250)[:2]:
                    draw.text((120, yb), line, fill="#" + C.INK, font=body_font)
                    yb += 36
                yb += 20
    draw.text((42, 862), "전남대학교 산업공학과 | MaintiQ Predict", fill="#" + C.MUTED, font=small_font)
    draw.text((1515, 862), str(idx), fill="#" + C.MUTED, font=small_font)
    return img


def build_preview() -> None:
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    for old_png in PREVIEW_DIR.glob("*.png"):
        old_png.unlink(missing_ok=True)
    thumbs = []
    for idx, spec in enumerate(SLIDES, start=1):
        img = preview_slide(spec, idx)
        out = PREVIEW_DIR / f"slide_{idx:02d}.png"
        img.save(out, quality=95)
        thumb = img.copy()
        thumb.thumbnail((360, 203))
        thumbs.append(thumb)
    cols = 4
    rows = math.ceil(len(thumbs) / cols)
    sheet = Image.new("RGB", (cols * 390, rows * 245 + 40), "white")
    draw = ImageDraw.Draw(sheet)
    for i, thumb in enumerate(thumbs):
        x = (i % cols) * 390 + 15
        y = (i // cols) * 245 + 15
        sheet.paste(thumb, (x, y))
        draw.text((x, y + 209), f"Slide {i + 1:02d}", fill="#555555", font=pil_font(16))
    sheet.save(CONTACT_SHEET, quality=95)


def write_notes() -> None:
    lines = ["# MaintiQ Predict 28장 최종발표 발표자 노트", ""]
    for idx, spec in enumerate(SLIDES, start=1):
        lines.extend(
            [
                f"## {idx}. {spec.title}",
                "",
                f"- 핵심 메시지: {spec.message}",
                f"- 발표 대본: {spec.talk}",
                f"- 예상 질문 대응: {spec.qa}",
                "",
            ]
        )
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8-sig")


def validate() -> None:
    prs = Presentation(PPTX_PATH)
    if len(prs.slides) != 28:
        raise RuntimeError(f"Expected 28 slides, got {len(prs.slides)}")
    combined = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                combined += shape.text + "\n"
    forbidden = [
        "실제 비용 절감 실증 완료",
        "실제 PLC/SCADA 배포 완료",
        "실제 회사 데이터 검증 완료",
        "실제 회사 데이터 성능 재검증 완료",
        "자동 정비 명령 실행 완료",
    ]
    for phrase in forbidden:
        if phrase in combined:
            raise RuntimeError(f"Forbidden phrase found: {phrase}")
    for pattern in [r"AIza[0-9A-Za-z_\-]{20,}", r"sk-[0-9A-Za-z_\-]{20,}"]:
        if re.search(pattern, combined):
            raise RuntimeError("API key-like pattern found in deck text")
    previews = list(PREVIEW_DIR.glob("slide_*.png"))
    if len(previews) != 28:
        raise RuntimeError(f"Expected 28 preview slides, got {len(previews)}")
    if not CONTACT_SHEET.exists() or CONTACT_SHEET.stat().st_size == 0:
        raise RuntimeError("Contact sheet missing")


def main() -> None:
    build_pptx()
    build_preview()
    write_notes()
    validate()
    print(f"PPTX: {PPTX_PATH}")
    print(f"Preview: {CONTACT_SHEET}")
    print(f"Notes: {NOTES_PATH}")


if __name__ == "__main__":
    main()
