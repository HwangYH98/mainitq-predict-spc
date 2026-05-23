"""Create thesis outline, references, speaker notes, and final presentation.

The deck is intentionally conservative and evidence-driven.  It uses the
project's existing metrics and charts, keeps the academic flow requested by the
user, and avoids unsupported claims about real factory deployment or actual
cost reduction.
"""

from __future__ import annotations

import csv
import json
import math
import textwrap
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = ROOT / "outputs"

PPTX_PATH = OUTPUTS / "final_presentation_maintiq_predict_academic_style.pptx"
NOTES_PATH = OUTPUTS / "final_presentation_academic_speaker_notes.md"
OUTLINE_PATH = OUTPUTS / "final_thesis_outline.md"
TOC_PATH = OUTPUTS / "final_thesis_table_of_contents.md"
REF_PATH = OUTPUTS / "research_references.md"
CANVA_OUTLINE_PATH = OUTPUTS / "final_presentation_canva_outline.md"
GENAI_EVIDENCE_PATH = OUTPUTS / "genai_report_evidence.md"
PROFESSOR_BRIEF_PATH = OUTPUTS / "professor_meeting_brief.md"
PREVIEW_DIR = OUTPUTS / "final_presentation_academic_preview"
CONTACT_SHEET = PREVIEW_DIR / "contact_sheet.png"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Malgun Gothic"


class C:
    GREEN = "003D12"
    GREEN2 = "008A3D"
    PALE = "EAF5EE"
    PALE2 = "F4FAF6"
    INK = "111827"
    MUTED = "596274"
    LINE = "C9D8CF"
    LIGHT = "F6F8FA"
    WHITE = "FFFFFF"
    BLUE = "1F6EA8"
    RED = "C0392B"
    AMBER = "D88A24"


def rgb(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def load_json(relative: str) -> dict:
    return json.loads((ROOT / relative).read_text(encoding="utf-8"))


def load_csv(relative: str) -> list[dict[str, str]]:
    with (ROOT / relative).open("r", encoding="utf-8-sig", newline="") as handle:
        return list(csv.DictReader(handle))


def as_float(value: str | float | int, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def fmt(value: str | float | int, digits: int = 3) -> str:
    return f"{as_float(value):.{digits}f}"


def pct(value: str | float | int, digits: int = 1) -> str:
    return f"{as_float(value) * 100:.{digits}f}%"


@dataclass
class Evidence:
    metrics: dict
    threshold: dict
    spc: dict
    model_strategy: list[dict[str, str]]
    spc_vs_ml: list[dict[str, str]]
    op_value: list[dict[str, str]]
    scania: dict
    public_bench: list[dict[str, str]]
    field: dict
    ai_report_context: dict
    ai_report_text: str


def load_evidence() -> Evidence:
    return Evidence(
        metrics=load_json("outputs/metrics.json"),
        threshold=load_json("outputs/threshold_summary.json"),
        spc=load_json("outputs/spc_summary.json"),
        model_strategy=load_csv("outputs/model_strategy_comparison.csv"),
        spc_vs_ml=load_csv("outputs/spc_vs_ml_comparison.csv"),
        op_value=load_csv("outputs/operational_value_simulation.csv"),
        scania=load_json("outputs/scania_official_cost_metrics.json"),
        public_bench=load_csv("outputs/public_industrial_validation_metrics.csv"),
        field=load_json("outputs/field_validation_report.json"),
        ai_report_context=load_json("outputs/ai_report_context.json"),
        ai_report_text=(ROOT / "outputs/ai_manager_report.md").read_text(encoding="utf-8", errors="replace"),
    )


def safe_path(path: Path) -> bool:
    return path.exists() and path.stat().st_size > 0


def blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(C.WHITE)
    return slide


def textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 12,
    color: str = C.INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_footer(slide, page_no: int) -> None:
    textbox(slide, "전남대학교 산업공학과 | MaintiQ Predict", 0.55, 7.08, 6.0, 0.2, 8, C.MUTED)
    textbox(slide, str(page_no), 12.45, 7.08, 0.35, 0.2, 8, C.MUTED, align=PP_ALIGN.RIGHT)


def title(slide, text: str, page_no: int) -> None:
    textbox(slide, text, 0.62, 0.42, 11.5, 0.45, 24, C.GREEN, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.02), Inches(11.95), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(C.GREEN2)
    line.line.color.rgb = rgb(C.GREEN2)
    add_footer(slide, page_no)


def shape_box(slide, x: float, y: float, w: float, h: float, fill: str = C.WHITE, line: str = C.LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(0.75)
    return shp


def section_label(slide, text: str, x: float, y: float, color: str = C.GREEN2) -> None:
    shape_box(slide, x, y, 1.4 + min(len(text), 12) * 0.05, 0.28, C.PALE, C.PALE)
    textbox(slide, text, x + 0.08, y + 0.045, 2.2, 0.16, 8, color, True)


def bullet_list(slide, items: Iterable[str], x: float, y: float, w: float, h: float, size: int = 12) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(C.INK)
        p.space_after = Pt(6)


def metric_card(slide, label: str, value: str, note: str, x: float, y: float, w: float = 2.55, h: float = 1.0, accent: str = C.GREEN2) -> None:
    shape_box(slide, x, y, w, h, C.WHITE, C.LINE)
    textbox(slide, label, x + 0.18, y + 0.16, w - 0.35, 0.18, 8, C.MUTED, True)
    textbox(slide, value, x + 0.18, y + 0.40, w - 0.35, 0.24, 16, accent, True)
    textbox(slide, note, x + 0.18, y + 0.70, w - 0.35, 0.18, 7, C.MUTED)


def mini_table(slide, rows: list[list[str]], x: float, y: float, col_w: list[float], row_h: float = 0.38, font_size: int = 8):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(sum(col_w)), Inches(row_h * len(rows)))
    table = table_shape.table
    for c, width in enumerate(col_w):
        table.columns[c].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            fill = C.GREEN if r == 0 else (C.PALE2 if r % 2 else C.WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill)
            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT
            para.font.size = Pt(font_size)
            para.font.bold = r == 0
            para.font.color.rgb = rgb(C.WHITE if r == 0 else C.INK)
    return table_shape


def add_picture(slide, path: Path, x: float, y: float, w: float, h: float | None = None) -> bool:
    if not safe_path(path):
        shape_box(slide, x, y, w, h or 1.5, C.LIGHT, C.LINE)
        textbox(slide, f"이미지 없음: {path.name}", x + 0.15, y + 0.2, w - 0.3, 0.4, 9, C.MUTED)
        return False
    if h is None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    return True


def add_bar_chart(slide, chart_title: str, categories: list[str], series: list[tuple[str, list[float]]], x: float, y: float, w: float, h: float) -> None:
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)


def add_flow(slide, items: list[tuple[str, str]], x: float, y: float, w: float, gap: float = 0.22) -> None:
    box_w = (w - gap * (len(items) - 1)) / len(items)
    for i, (head, body) in enumerate(items):
        bx = x + i * (box_w + gap)
        shape_box(slide, bx, y, box_w, 1.0, C.PALE2, C.LINE)
        textbox(slide, head, bx + 0.12, y + 0.18, box_w - 0.24, 0.18, 9, C.GREEN, True, PP_ALIGN.CENTER)
        textbox(slide, body, bx + 0.12, y + 0.52, box_w - 0.24, 0.2, 8, C.INK, align=PP_ALIGN.CENTER)
        if i < len(items) - 1:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(bx + box_w),
                Inches(y + 0.5),
                Inches(bx + box_w + gap),
                Inches(y + 0.5),
            )
            conn.line.color.rgb = rgb(C.GREEN2)
            conn.line.width = Pt(1.2)


def group_public(rows: list[dict[str, str]]) -> dict[str, list[dict[str, str]]]:
    out: dict[str, list[dict[str, str]]] = {}
    for row in rows:
        out.setdefault(row["dataset_id"], []).append(row)
    return out


def best_public_row(rows: list[dict[str, str]]) -> dict[str, str]:
    for strategy in ("xgboost_tuned_threshold", "xgboost_default", "ml_spc_combined", "rule_based_threshold"):
        for row in rows:
            if row.get("strategy_id") == strategy:
                return row
    return rows[0]


def ai_report_summary(ev: Evidence) -> dict[str, str]:
    """Return slide/paper-ready GenAI report evidence without exposing secrets."""
    ctx = ev.ai_report_context
    row = ctx.get("row", {})
    future = ctx.get("future_prediction", {})
    factors = ctx.get("top_shap_factors", [])
    factor_names = []
    for factor in factors[:3]:
        name = str(factor.get("feature", "")).replace("_", " ")
        if name:
            factor_names.append(name)
    return {
        "mode": str(ctx.get("report_generation_mode", "unknown")),
        "generated_at": str(ctx.get("generated_at", "-")),
        "probability": fmt(row.get("xgboost_probability", 0), 6),
        "threshold": fmt(row.get("selected_threshold", 0), 2),
        "risk_status": str(row.get("risk_status", "-")),
        "udi": str(row.get("UDI", "-")),
        "timestamp": str(row.get("simulated_timestamp", "-")),
        "future_probability": fmt(future.get("predicted_future_deviation_probability_h10", 0), 6),
        "future_deviation": str(future.get("predicted_future_deviation_h10", "-")),
        "factors": ", ".join(factor_names) if factor_names else "torque, rotational speed",
        "guardrail": str(ctx.get("guardrail", "관리자 참고용이며 자동 정비 명령이 아님")),
    }


def write_references() -> None:
    REF_PATH.write_text(
        """# MaintiQ Predict 연구 참고문헌

아래 참고문헌은 논문 본문과 최종발표 자료에 사용할 수 있도록 주제별로 정리한 목록이다. 최종 논문 제출 시 학교 양식에 맞춰 번호식 또는 APA 형식으로 통일한다.

## 1. 예지보전, CBM, 산업공학 배경

1. Jardine, A. K. S., Lin, D., & Banjevic, D. (2006). A review on machinery diagnostics and prognostics implementing condition-based maintenance. *Mechanical Systems and Signal Processing, 20*(7), 1483-1510.
2. Carvalho, T. P., Soares, F. A. A. M. N., Vita, R., Francisco, R. da P., Basto, J. P., & Alcalá, S. G. S. (2019). A systematic literature review of machine learning methods applied to predictive maintenance. *Computers & Industrial Engineering, 137*, 106024.
3. Montgomery, D. C. (2019). *Introduction to Statistical Quality Control* (8th ed.). Wiley.
4. Elkan, C. (2001). The foundations of cost-sensitive learning. *Proceedings of IJCAI*, 973-978.

## 2. 공개 산업 데이터셋

5. UCI Machine Learning Repository. (2020). *AI4I 2020 Predictive Maintenance Dataset*.
6. Matzka, S. (2020). Explainable Artificial Intelligence for Predictive Maintenance Applications. *AI4I 2020*.
7. Kharazian, Z., Lindgren, T., Magnússon, S., Steinert, O., & Andersson Reyna, O. (2024). *SCANIA Component X Dataset: A Real-World Multivariate Time Series Dataset for Predictive Maintenance*. Researchdata.se.
8. Kharazian, Z., Lindgren, T., Magnússon, S., Steinert, O., & Andersson Reyna, O. (2024). SCANIA Component X Dataset: A Real-World Multivariate Time Series Dataset for Predictive Maintenance. *arXiv:2401.15199*.
9. UCI Machine Learning Repository. (2023). *MetroPT-3 Dataset*.
10. Saxena, A., & Goebel, K. (2008). *Turbofan Engine Degradation Simulation Data Set*. NASA Ames Prognostics Data Repository.
11. NASA. (2023). *IMS Bearings*.

## 3. 모델링, 불균형 처리, 설명 가능성

12. Chen, T., & Guestrin, C. (2016). XGBoost: A scalable tree boosting system. *KDD*, 785-794.
13. Chawla, N. V., Bowyer, K. W., Hall, L. O., & Kegelmeyer, W. P. (2002). SMOTE: Synthetic Minority Over-sampling Technique. *JAIR, 16*, 321-357.
14. Lundberg, S. M., & Lee, S.-I. (2017). A unified approach to interpreting model predictions. *NeurIPS, 30*.
15. scikit-learn developers. Probability calibration. *scikit-learn User Guide*.

## 4. GenAI 리포트 및 설치형 앱 구현

16. Google AI for Developers. Gemini API: models.generateContent.
17. OpenAI. Responses API reference.
18. PyInstaller Development Team. Using PyInstaller.

## 발표용 축약 References

- Jardine et al. (2006); Carvalho et al. (2019); Montgomery (2019)
- AI4I 2020; Matzka (2020); SCANIA Component X; MetroPT-3; NASA C-MAPSS; IMS Bearings
- XGBoost; SMOTE; SHAP; cost-sensitive learning
- Gemini generateContent API; OpenAI Responses API; PyInstaller
""",
        encoding="utf-8",
    )


def write_thesis_docs(ev: Evidence) -> None:
    xgb = ev.metrics["models"]["xgboost"]
    tuned = ev.threshold["selected_metrics"]
    scania_best = next(row for row in ev.scania["metrics"] if row["strategy_id"] == "xgboost_cost_optimized")
    ai = ai_report_summary(ev)

    OUTLINE_PATH.write_text(
        f"""# MaintiQ Predict 논문 개요

## 연구 제목 후보
**ML 예측, Predictive SPC, GenAI 리포트 및 승인형 작업지시를 결합한 스마트 제조 예지보전 운영 시스템 구현**

## 초록 초안
본 연구는 스마트 제조 환경에서 설비 센서 데이터를 기반으로 고장 위험을 예측하고, 예측 결과를 관리도 기반 모니터링, GenAI 관리자 리포트, 승인형 작업지시 흐름으로 연결하는 예지보전 운영 MVP를 구현한다. AI4I 2020 데이터셋을 기준으로 Logistic Regression, XGBoost, SMOTE 적용 모델, threshold tuning을 비교했으며, XGBoost는 PR-AUC {xgb['pr_auc']:.4f}를 보였다. F1 기준 threshold tuning에서는 threshold {ev.threshold['selected_threshold']:.2f}, precision {tuned['precision']:.4f}, recall {tuned['recall']:.4f}, F1-score {tuned['f1_score']:.4f}를 얻었다. Gemini 기반 GenAI 리포트 검증에서는 `report_generation_mode={ai['mode']}`로 관리자 참고 리포트를 생성했고, 예측 확률 {ai['probability']}, 임계치 {ai['threshold']}, 상태 {ai['risk_status']}를 요약했다. 또한 SPC-only, ML threshold, ML+SPC 정책을 비교하고 false alarm과 missed failure 기반 normalized cost simulation을 수행했다. 공개 산업 데이터인 SCANIA Component X에 대해서는 official class 0~4 cost matrix를 적용해 rule baseline 대비 official cost metric {scania_best['cost_improvement_vs_rule'] * 100:.2f}% 개선 가능성을 확인했다. 단, 본 연구는 실제 PLC/SCADA 운영망 배포나 실제 회사 비용 절감 실증이 아니라 공개 데이터와 로컬 MVP 기반 재현 가능 검증에 초점을 둔다.

## 연구 배경
- 사후보전은 고장 발생 후 대응하므로 downtime과 정비 비용이 커질 수 있다.
- 예방보전은 일정 기반 정비로 과잉 정비와 누락 위험이 존재한다.
- 예지보전은 센서 데이터로 고장 가능성을 사전에 추정하지만, 예측 결과가 현장 의사결정으로 연결되지 않으면 운영 가치가 제한된다.

## 연구 목적
1. AI4I 2020 기반 고장 예측 모델과 threshold 정책을 구현한다.
2. Predictive SPC를 통해 위험 확률의 시간 흐름과 관리 한계를 함께 제시한다.
3. GenAI 리포트와 승인형 작업지시 workflow로 예측 결과를 운영 의사결정 흐름에 연결한다.
4. SCANIA Component X와 공개 benchmark를 통해 논문에서 방어 가능한 공개 데이터 검증 근거를 확보한다.

## GenAI 리포트 검증 요약
- 생성 방식: `{ai['mode']}`
- 입력 context: AI4I UDI {ai['udi']}, 시뮬레이션 시각 {ai['timestamp']}
- 위험 판단: 예측 확률 {ai['probability']}, 임계치 {ai['threshold']}, 상태 {ai['risk_status']}
- 주요 위험 요인: {ai['factors']}
- 관리자 조치 방향: 현장 확인, 토크/회전속도 관련 부품 점검, 승인형 작업지시 검토
- 한계: AI4I 기반 offline simulation이며 실제 현장 센서 feed 또는 자동 정비 명령이 아니다.

## 차별성
- 단일 분류 모델 성능만 제시하지 않고, 예측 → SPC → 리포트 → 작업지시 → 감사 이력까지 연결한다.
- Lite/Full 설치형 데스크톱 앱과 Admin 검증 콘솔을 분리해 사용자 흐름과 연구 검증 흐름을 구분한다.
- 공개 산업 데이터의 official cost metric과 normalized cost simulation을 함께 사용해 비용 민감 의사결정 가능성을 평가한다.

## 주장 경계
- 가능: 공개 데이터 기준 모델 비교, threshold trade-off, SPC-only 대비 ML+SPC 정책 비교, SCANIA official cost metric 기준 개선 가능성.
- 금지: 실제 PLC/SCADA 운영망 배포 완료, 실제 회사 데이터 성능 검증 완료, 실제 원화 비용 절감 실증 완료, 자동 정비를 시스템이 직접 수행했다는 표현.
""",
        encoding="utf-8",
    )


def write_genai_evidence_docs(ev: Evidence) -> None:
    ai = ai_report_summary(ev)
    xgb = ev.metrics["models"]["xgboost"]
    tuned = ev.threshold["selected_metrics"]
    scania_best = next(row for row in ev.scania["metrics"] if row["strategy_id"] == "xgboost_cost_optimized")

    GENAI_EVIDENCE_PATH.write_text(
        f"""# Gemini AI 리포트 검증 근거

## 생성 결과 요약

| 항목 | 값 |
|---|---|
| report_generation_mode | `{ai['mode']}` |
| 생성 시각 | {ai['generated_at']} |
| 입력 row | UDI {ai['udi']} / {ai['timestamp']} |
| 예측 고장 확률 | {ai['probability']} |
| 위험 판정 기준 | {ai['threshold']} |
| 위험 상태 | {ai['risk_status']} |
| 주요 위험 요인 | {ai['factors']} |
| 향후 10-step 이탈 확률 | {ai['future_probability']} |
| 향후 10-step 이탈 예측 | {ai['future_deviation']} |

## 논문/PPT에 넣을 축약 문장

Gemini generateContent API를 통해 AI4I 기반 위험 context를 관리자 참고 리포트로 변환하였다. 리포트는 UDI {ai['udi']} 설비 row에 대해 예측 고장 확률 {ai['probability']}, 위험 판정 기준 {ai['threshold']}, 상태 {ai['risk_status']}를 제시하고, 주요 위험 요인으로 {ai['factors']}를 요약하였다. 생성된 리포트는 현장 확인과 부품 점검을 권고하지만, 자동 정비 명령이 아니라 승인형 작업지시 검토를 위한 참고 자료로 제한된다.

## 발표 시 주의 문장

- API key는 세션에서만 사용하며 파일이나 Git 기록에 저장하지 않는다.
- 본 리포트는 AI4I 기반 offline simulation 결과이며 실제 현장 센서 feed가 아니다.
- 최종 조치는 현장 담당자의 확인과 승인 후 진행해야 한다.

## 원문 산출물

- `outputs/ai_report_context.json`
- `outputs/ai_manager_report.md`
""",
        encoding="utf-8",
    )

    PROFESSOR_BRIEF_PATH.write_text(
        f"""# 교수님 면담용 요약 자료

## 1. 연구 주제
**MaintiQ Predict: ML 예측, Predictive SPC, GenAI 리포트, 승인형 작업지시를 결합한 스마트 제조 예지보전 운영 시스템**

## 2. 현재 구현 범위
- AI4I 2020 기반 고장 예측 모델과 threshold tuning 구현
- Predictive SPC와 ML 위험 확률 모니터링 구현
- Gemini 기반 관리자 참고 리포트 생성 검증
- 승인형 작업지시 workflow와 운영 이력 저장 구현
- SCANIA Component X official cost metric 검증 및 public benchmark 확장
- Full/Lite Windows 데스크톱 앱과 Streamlit Admin 콘솔 구현

## 3. 핵심 결과
- XGBoost PR-AUC: {xgb['pr_auc']:.4f}
- threshold tuning 후 F1-score: {tuned['f1_score']:.4f}
- 위험 판정 기준: {ev.threshold['selected_threshold']:.2f}
- Gemini 리포트: `{ai['mode']}`로 생성, 예측 확률 {ai['probability']}, 상태 {ai['risk_status']}
- SCANIA official cost metric: rule baseline 대비 {pct(scania_best['cost_improvement_vs_rule'], 2)} 개선 가능성

## 4. 교수님께 확인받을 내용
1. 논문 제목을 제품명 중심으로 둘지, 연구 방법 중심으로 둘지
2. SCANIA 결과를 “공개 산업 benchmark official cost metric 개선”으로 표현하는 것이 적절한지
3. Gemini AI 리포트를 본문 실험 결과에 넣을지, 시스템 구현 장에 넣을지
4. 실제 회사 데이터 실증이 없다는 한계를 결론에서 어느 정도 강조할지

## 5. 금지할 주장
- 실제 PLC/SCADA 운영망 배포 완료
- 실제 회사 데이터 성능 검증 완료
- 실제 원화 비용 절감 실증 완료
- 자동 정비를 시스템이 직접 수행했다는 표현
""",
        encoding="utf-8",
    )

def write_toc() -> None:
    TOC_PATH.write_text(
        """# MaintiQ Predict 논문 목차

## 1. 서론
1.1 연구 배경  
1.2 연구 목적  
1.3 연구 범위 및 제한  
1.4 연구 차별성  
1.5 논문 구성  

## 2. 이론적 배경 및 선행연구
2.1 사후보전, 예방보전, 예지보전  
2.2 Condition-Based Maintenance와 Predictive Maintenance  
2.3 Statistical Process Control와 관리도  
2.4 ML 기반 고장 예측과 불균형 데이터  
2.5 설명 가능 AI와 GenAI 리포트  
2.6 비용 민감 학습과 공개 산업 benchmark  

## 3. 연구 방법
3.1 전체 연구 절차  
3.2 데이터 구성과 전처리  
3.3 고장 예측 모델과 threshold tuning  
3.4 Predictive SPC 구성  
3.5 GenAI 리포트와 승인형 작업지시 workflow  
3.6 실험 설계와 평가지표  

## 4. 시스템 구현
4.1 MaintiQ Predict 데스크톱 앱  
4.2 빠른 점검 모드와 정밀 분석 모드  
4.3 Streamlit Admin 검증 콘솔  
4.4 감사 로그와 작업지시 이력  
4.5 GenAI 관리자 리포트 생성 구조  
4.6 설치형 배포 구조  

## 5. 실험 및 검증
5.1 AI4I baseline 결과  
5.2 SMOTE 및 threshold tuning 비교  
5.3 SPC-only vs ML+SPC 비교  
5.4 운영 cost simulation  
5.5 Gemini API 기반 관리자 리포트 생성 검증  
5.6 SCANIA official cost metric 검증  
5.7 MetroPT-3, C-MAPSS, IMS/FEMTO benchmark 확장  
5.8 실제 회사 데이터 실증 준비성과 한계  

## 6. 결론 및 향후 연구
6.1 연구 결과 요약  
6.2 산업공학적 의의  
6.3 연구 한계  
6.4 향후 연구  

## 7. 참고문헌
예지보전, SPC, XGBoost, SMOTE, SHAP, cost-sensitive learning, AI4I, SCANIA, MetroPT-3, NASA C-MAPSS, IMS Bearings, GenAI API, PyInstaller 관련 문헌을 포함한다.
""",
        encoding="utf-8",
    )


def slide_cover(prs: Presentation) -> None:
    slide = blank(prs)
    textbox(slide, "MaintiQ Predict", 0.75, 0.75, 8.0, 0.6, 30, C.GREEN, True)
    textbox(slide, "ML·SPC·GenAI를 결합한 AI 예지보전 운영 시스템", 0.78, 1.42, 10.5, 0.38, 17, C.INK, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(8.6), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(C.GREEN2)
    line.line.color.rgb = rgb(C.GREEN2)
    textbox(slide, "학부 캡스톤 최종발표", 0.82, 2.25, 5.5, 0.28, 12, C.MUTED)
    textbox(slide, "예측 모델 성능 비교에서 운영 리포트와 승인형 작업지시까지 연결하는 로컬 제품형 MVP", 0.82, 2.75, 9.8, 0.3, 13, C.INK)

    # Large editorial evidence strip instead of generic KPI cards.
    shape_box(slide, 0.82, 4.45, 11.5, 1.35, C.PALE2, C.LINE)
    textbox(slide, "핵심 메시지", 1.05, 4.72, 1.8, 0.25, 11, C.GREEN, True)
    textbox(slide, "고장 확률 예측 결과를 관리도, GenAI 리포트, 작업지시 승인 이력으로 연결했다.", 2.55, 4.72, 8.6, 0.32, 15, C.INK, True)
    textbox(slide, "단, 실제 공장 배포나 실제 비용 절감 실증이 아니라 공개 데이터 기반 검증과 현장 실증 준비 단계로 한정한다.", 1.05, 5.24, 10.3, 0.22, 9, C.MUTED)
    add_footer(slide, 1)


def slide_agenda(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "목차", 2)
    items = [
        ("1", "Motivation", "문제 배경과 연구 필요성"),
        ("2", "Contributions", "본 연구의 구현 및 검증 기여"),
        ("3", "참고 문헌 Review", "예지보전, ML, SPC, 비용 민감 학습"),
        ("4", "Main Part", "시스템 구조, 데이터, 모델, workflow"),
        ("5", "실험 설계", "비교군과 평가지표"),
        ("6", "실험 결과", "AI4I, SCANIA, public benchmark"),
        ("7", "결론 및 References", "주장 경계와 향후 연구"),
    ]
    for i, (no, head, body) in enumerate(items):
        y = 1.35 + i * 0.68
        textbox(slide, no, 1.0, y, 0.4, 0.25, 15, C.GREEN, True)
        textbox(slide, head, 1.55, y, 3.2, 0.25, 15, C.INK, True)
        textbox(slide, body, 4.85, y + 0.04, 5.2, 0.22, 10, C.MUTED)


def slide_motivation(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "Motivation", 3)
    textbox(slide, "예측 결과가 현장 의사결정으로 연결되지 않으면 운영 가치는 제한된다.", 0.78, 1.25, 11.0, 0.35, 17, C.INK, True)
    add_flow(
        slide,
        [
            ("센서 데이터", "공정 상태"),
            ("고장 확률", "ML 예측"),
            ("위험 추세", "SPC"),
            ("관리자 리포트", "GenAI"),
            ("작업지시", "승인 이력"),
        ],
        0.78,
        2.0,
        11.5,
    )
    bullet_list(
        slide,
        [
            "기존 threshold 또는 단일 SPC만으로는 복합 센서 조합과 고장 확률 변화를 충분히 설명하기 어렵다.",
            "모델 성능이 좋아도 위험 원인, 조치 근거, 작업자 결정 이력이 없으면 운영 도입성이 낮다.",
            "따라서 본 연구는 예측 모델 자체보다 예측 결과를 운영 workflow로 연결하는 데 초점을 둔다.",
        ],
        0.95,
        3.65,
        10.8,
        1.8,
        13,
    )
    metric_card(slide, "연구 위치", "로컬 제품형 MVP", "공개 데이터 검증 기반", 0.95, 6.05, 3.0)
    metric_card(slide, "핵심 제약", "현장 실증 전", "회사 로그 필요", 4.25, 6.05, 3.0, accent=C.AMBER)
    metric_card(slide, "운영 방식", "승인형 작업지시", "자동 명령 아님", 7.55, 6.05, 3.0, accent=C.BLUE)


def slide_contributions(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "Contributions", 4)
    rows = [
        ["기여", "구현 내용", "검증 근거"],
        ["예측", "AI4I 기반 Logistic Regression, XGBoost, SMOTE, threshold tuning", "PR-AUC, F1, precision/recall"],
        ["모니터링", "ML probability와 Predictive SPC 결합", "SPC-only vs ML+SPC 비교"],
        ["운영", "GenAI 리포트와 승인형 작업지시 workflow", "SQLite/CSV 이력, report history"],
        ["공개 검증", "SCANIA official cost metric 및 public benchmark adapter", "official cost, lead-time, simulated cost"],
    ]
    mini_table(slide, rows, 0.65, 1.35, [2.1, 5.5, 3.7], 0.56, 8.5)
    textbox(slide, "Something New", 0.78, 5.35, 2.6, 0.3, 18, C.GREEN, True)
    textbox(slide, "단일 예측 모델이 아니라 예측, SPC, GenAI 설명, 승인형 작업지시, 공개 benchmark 검증을 하나의 실행 가능한 시스템으로 통합했다.", 0.82, 5.85, 10.9, 0.45, 15, C.INK, True)


def slide_review(prs: Presentation, page: int, review_no: int, heading: str, rows: list[list[str]], takeaway: str) -> None:
    slide = blank(prs)
    title(slide, f"Review {review_no}: {heading}", page)
    mini_table(slide, rows, 0.7, 1.35, [2.6, 4.15, 4.65], 0.5, 8)
    textbox(slide, "본 연구 반영점", 0.82, 5.15, 2.0, 0.28, 15, C.GREEN, True)
    shape_box(slide, 0.82, 5.62, 11.0, 0.72, C.PALE2, C.LINE)
    textbox(slide, takeaway, 1.08, 5.83, 10.3, 0.24, 12, C.INK, True)


def slide_system(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "Main Part: Proposed System", 8)
    add_flow(
        slide,
        [
            ("1. CSV 입력", "센서 row"),
            ("2. 전처리", "컬럼/단위"),
            ("3. 예측", "위험 확률"),
            ("4. SPC", "관리도"),
            ("5. AI 리포트", "관리자 요약"),
            ("6. 작업지시", "승인 이력"),
        ],
        0.55,
        1.45,
        12.1,
    )
    add_picture(slide, OUTPUTS / "maintiq_predict_screenshot.png", 0.78, 3.05, 7.0, 3.25)
    bullet_list(
        slide,
        [
            "사용자 앱은 제품형 데스크톱 UI로 구성하고 연구/검증 화면은 Admin 콘솔로 분리했다.",
            "Lite는 빠른 점검 모드, Full은 정밀 분석 모드로 구분한다.",
            "API key와 원본 회사 데이터는 파일로 저장하지 않는 것을 기본 정책으로 둔다.",
        ],
        8.05,
        3.25,
        4.2,
        2.2,
        11,
    )


def slide_data(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "Main Part: Data & Preprocessing", 9)
    rows = [
        ["구분", "처리 내용"],
        ["기준 데이터", "AI4I 2020 Predictive Maintenance Dataset"],
        ["목표 변수", "Machine failure binary target"],
        ["제거 컬럼", "UDI, Product ID, TWF/HDF/PWF/OSF/RNF leakage columns"],
        ["범주형 처리", "Type one-hot encoding"],
        ["회사 CSV", "컬럼 alias mapping, Celsius/Kelvin 변환, 결측/이상값 진단"],
    ]
    mini_table(slide, rows, 0.75, 1.35, [2.6, 8.3], 0.5, 9)
    add_flow(
        slide,
        [
            ("업로드", "CSV"),
            ("매핑", "필수 컬럼"),
            ("품질 진단", "결측/단위"),
            ("예측", "확률/우선순위"),
            ("저장", "결과 CSV"),
        ],
        1.0,
        5.15,
        10.8,
        0.25,
    )


def slide_model(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "Main Part: Model and Threshold Policy", 10)
    xgb = ev.metrics["models"]["xgboost"]
    logit = ev.metrics["models"]["logistic_regression"]
    tuned = ev.threshold["selected_metrics"]
    add_bar_chart(
        slide,
        "Baseline metric comparison",
        ["Precision", "Recall", "F1", "ROC-AUC", "PR-AUC"],
        [
            ("Logistic", [logit[k] for k in ["precision", "recall", "f1_score", "roc_auc", "pr_auc"]]),
            ("XGBoost", [xgb[k] for k in ["precision", "recall", "f1_score", "roc_auc", "pr_auc"]]),
        ],
        0.75,
        1.35,
        6.8,
        3.45,
    )
    metric_card(slide, "선택 모델", "XGBoost", f"PR-AUC {xgb['pr_auc']:.4f}", 8.05, 1.45, 3.5)
    metric_card(slide, "위험 판정 기준", f"{ev.threshold['selected_threshold']:.2f}", "F1 기준 threshold", 8.05, 2.72, 3.5, accent=C.BLUE)
    metric_card(slide, "튜닝 후 F1", f"{tuned['f1_score']:.4f}", f"P {tuned['precision']:.3f} / R {tuned['recall']:.3f}", 8.05, 3.99, 3.5, accent=C.GREEN2)
    textbox(slide, "불균형 데이터에서는 accuracy보다 PR-AUC, recall, false alarm/missed failure trade-off를 함께 봐야 한다.", 0.9, 5.55, 10.7, 0.35, 12, C.MUTED)


def slide_spc(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "Main Part: Predictive SPC", 11)
    add_picture(slide, OUTPUTS / "spc_risk_chart.png", 0.65, 1.25, 6.2, 3.0)
    add_picture(slide, OUTPUTS / "spc_control_chart.png", 0.65, 4.55, 6.2, 1.7)
    metric_card(slide, "관측 row", str(ev.spc["total_rows"]), "AI4I test prediction", 7.25, 1.45, 3.8)
    metric_card(slide, "고위험 row", str(ev.spc["high_risk_count"]), "threshold 초과", 7.25, 2.72, 3.8, accent=C.RED)
    metric_card(slide, "SPC risk alert", str(ev.spc["spc_risk_alert_count"]), "위험 확률 관리도 기준", 7.25, 3.99, 3.8, accent=C.AMBER)
    textbox(slide, "SPC는 단일 센서 관리도만이 아니라 ML 위험 확률의 시간 흐름을 함께 감시하는 방식으로 확장했다.", 7.35, 5.35, 4.2, 0.55, 11, C.INK, True)


def slide_workflow(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "Main Part: GenAI Report and Work-order Workflow", 12)
    ai = ai_report_summary(ev)
    add_flow(
        slide,
        [
            ("위험 row", "상위 설비"),
            ("리포트", "요약/원인"),
            ("초안", "작업지시"),
            ("결정", "승인/검토/반려"),
            ("이력", "감사 로그"),
        ],
        0.85,
        1.45,
        11.0,
    )
    textbox(
        slide,
        f"Gemini 실호출 결과: {ai['mode']} | 위험 상태 {ai['risk_status']}",
        0.85,
        2.62,
        10.8,
        0.25,
        10,
        C.GREEN,
        True,
    )
    rows = [
        ["리포트 항목", "실제 생성 결과", "발표/논문 반영"],
        ["연동 방식", ai["mode"], "Gemini API 실호출 증거"],
        ["위험 판단", f"p={ai['probability']} > 기준 {ai['threshold']} / {ai['risk_status']}", "고위험 row 요약"],
        ["주요 요인", ai["factors"], "SHAP 기반 설명 근거"],
        ["관리자 조치", "현장 확인, 주요 부품 점검, 승인형 작업지시 검토", "자동 정비 명령 아님"],
    ]
    mini_table(slide, rows, 0.7, 3.05, [2.0, 4.55, 4.65], 0.46, 7.8)
    metric_card(slide, "입력 row", f"UDI {ai['udi']}", ai["timestamp"], 0.85, 5.72, 2.75, 0.9)
    metric_card(slide, "예측 확률", ai["probability"], f"기준 {ai['threshold']}", 3.85, 5.72, 2.75, 0.9, accent=C.RED)
    metric_card(slide, "향후 이탈 확률", ai["future_probability"], f"10-step 예측 {ai['future_deviation']}", 6.85, 5.72, 2.75, 0.9, accent=C.AMBER)
    textbox(slide, "API key는 세션에서만 사용하며 파일·README·Git 기록에 저장하지 않는다.", 9.85, 5.86, 2.2, 0.45, 8, C.MUTED)


def slide_experiment_design(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "실험 설계", 13)
    rows = [
        ["실험", "비교군", "평가지표"],
        ["AI4I baseline", "Logistic Regression, XGBoost", "precision, recall, F1, ROC-AUC, PR-AUC"],
        ["불균형/threshold", "SMOTE, tuned threshold", "F1, alert count, FP/FN"],
        ["SPC 비교", "SPC-only, ML threshold, ML+SPC", "F1, recall, false alarm, missed failure"],
        ["운영 가치", "no alert, rule, tuned ML, ML+SPC", "normalized cost"],
        ["공개 산업 데이터", "SCANIA, MetroPT-3, C-MAPSS, IMS/FEMTO", "official cost, lead-time, RUL"],
    ]
    mini_table(slide, rows, 0.55, 1.25, [2.35, 4.1, 4.7], 0.48, 8.2)
    textbox(slide, "실험 설계 원칙", 0.8, 5.35, 2.1, 0.28, 15, C.GREEN, True)
    bullet_list(
        slide,
        [
            "같은 split과 같은 threshold 정책에서 비교해 모델·정책 효과를 분리한다.",
            "실제 비용 절감은 주장하지 않고, 공개 데이터에서는 official metric 또는 normalized simulation으로 제한한다.",
        ],
        0.95,
        5.78,
        10.8,
        0.8,
        12,
    )


def slide_ai4i(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "실험 결과: AI4I Baseline", 14)
    xgb = ev.metrics["models"]["xgboost"]
    logit = ev.metrics["models"]["logistic_regression"]
    rows = [
        ["모델", "Precision", "Recall", "F1", "ROC-AUC", "PR-AUC"],
        ["Logistic Regression", fmt(logit["precision"], 4), fmt(logit["recall"], 4), fmt(logit["f1_score"], 4), fmt(logit["roc_auc"], 4), fmt(logit["pr_auc"], 4)],
        ["XGBoost", fmt(xgb["precision"], 4), fmt(xgb["recall"], 4), fmt(xgb["f1_score"], 4), fmt(xgb["roc_auc"], 4), fmt(xgb["pr_auc"], 4)],
    ]
    mini_table(slide, rows, 0.6, 1.35, [3.1, 1.35, 1.35, 1.35, 1.35, 1.35], 0.52, 8.5)
    add_picture(slide, OUTPUTS / "pr_curve.png", 3.35, 3.25, 5.8, 2.8)
    metric_card(slide, "Best baseline", "XGBoost", f"PR-AUC {xgb['pr_auc']:.4f}", 9.55, 3.35, 2.65)
    metric_card(slide, "Test failures", str(ev.threshold["test_failures"]), "불균형 target", 9.55, 4.62, 2.65, accent=C.AMBER)


def slide_strategy(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "실험 결과: SMOTE / Threshold 비교", 15)
    display = []
    preferred = [
        "xgboost_default_tuned_threshold",
        "xgboost_smote_tuned_threshold",
        "xgboost_default",
        "logistic_regression_default",
    ]
    for sid in preferred:
        row = next((r for r in ev.model_strategy if r["strategy_id"] == sid), None)
        if row:
            display.append(row)
    rows = [["전략", "Threshold", "Precision", "Recall", "F1", "PR-AUC"]]
    rows += [[r["display_name"], r["threshold"], r["precision"], r["recall"], r["f1_score"], r["pr_auc"]] for r in display]
    mini_table(slide, rows, 0.55, 1.35, [3.6, 1.05, 1.25, 1.25, 1.25, 1.25], 0.5, 8)
    bullet_list(
        slide,
        [
            "본 split에서는 XGBoost + tuned threshold가 F1 기준 가장 안정적인 결과를 보였다.",
            "SMOTE는 recall을 보강할 수 있지만 false alarm과 precision trade-off가 발생할 수 있어 자동 채택 대상이 아니다.",
            "따라서 운영 정책은 모델 자체보다 threshold와 cost policy를 함께 조정해야 한다.",
        ],
        0.85,
        4.65,
        10.8,
        1.25,
        12,
    )


def slide_spc_compare(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "실험 결과: SPC-only vs ML+SPC", 16)
    rows = [["정책", "Precision", "Recall", "F1", "Alert count"]]
    rows += [[r["display_name"], r["precision"], r["recall"], r["f1_score"], r["alert_count"]] for r in ev.spc_vs_ml]
    mini_table(slide, rows, 0.75, 1.25, [4.0, 1.35, 1.35, 1.35, 1.4], 0.52, 8.5)
    add_bar_chart(
        slide,
        "F1 / Recall by alert policy",
        [r["display_name"].replace(" torque control limit", "").replace(" selected threshold", "") for r in ev.spc_vs_ml],
        [
            ("F1", [as_float(r["f1_score"]) for r in ev.spc_vs_ml]),
            ("Recall", [as_float(r["recall"]) for r in ev.spc_vs_ml]),
        ],
        1.4,
        4.0,
        9.5,
        2.2,
    )


def slide_cost(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "실험 결과: Operational Cost Simulation", 17)
    rows = [["정책", "Alerts", "False alarm", "Missed failure", "Normalized cost"]]
    rows += [
        [r["display_name"], r["alert_count"], r["false_alarm_count"], r["missed_failure_count"], r["normalized_operating_cost"]]
        for r in ev.op_value
        if r["scenario_id"] == "conservative"
    ]
    mini_table(slide, rows, 0.55, 1.25, [3.8, 1.2, 1.45, 1.65, 1.7], 0.4, 8)
    add_picture(slide, OUTPUTS / "operational_value_simulation.png", 2.35, 4.15, 8.1, 2.1)
    textbox(slide, "주의: 이 값은 실제 원화 비용 절감 실증이 아니라 false alarm/missed failure 가중치 기반 상대 cost simulation이다.", 0.75, 6.55, 11.0, 0.25, 9, C.MUTED)


def slide_scania(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "실험 결과: SCANIA Official Cost Metric", 18)
    best = next(row for row in ev.scania["metrics"] if row["strategy_id"] == "xgboost_cost_optimized")
    metric_card(slide, "공개 실제 산업 데이터", "SCANIA", "Component X fleet dataset", 0.85, 1.35, 3.35)
    metric_card(slide, "평가 기준", "Official cost", "class 0~4 cost matrix", 4.45, 1.35, 3.35, accent=C.BLUE)
    metric_card(slide, "Rule 대비 개선", pct(best["cost_improvement_vs_rule"], 2), "official cost metric", 8.05, 1.35, 3.35, accent=C.GREEN2)
    add_picture(slide, OUTPUTS / "scania_official_cost_comparison.png", 0.85, 2.9, 5.35, 3.1)
    add_picture(slide, OUTPUTS / "scania_official_confusion_matrix.png", 6.65, 2.9, 4.9, 3.1)
    textbox(slide, "표현 범위: 실제 회사 원가 절감이 아니라 SCANIA 공개 benchmark의 official cost metric 기준 개선이다.", 0.9, 6.35, 10.8, 0.25, 9, C.MUTED)


def slide_public_bench(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "실험 결과: Public Benchmark Extension", 19)
    grouped = group_public(ev.public_bench)
    rows = [["Dataset", "분석 범위", "대표 전략", "핵심 지표"]]
    for dataset_id, group in list(grouped.items())[:5]:
        row = best_public_row(group)
        rows.append([
            dataset_id.upper(),
            row.get("label_scope", "-")[:28],
            row.get("display_name", "-")[:32],
            f"F1 {as_float(row.get('f1_score')):.3f} / PR-AUC {as_float(row.get('pr_auc')):.3f}",
        ])
    mini_table(slide, rows, 0.55, 1.35, [1.55, 3.35, 3.45, 2.4], 0.46, 7.8)
    bullet_list(
        slide,
        [
            "MetroPT-3는 압축기 이상 구간 기반 lead-time 분석에 적합하다.",
            "C-MAPSS는 RUL 및 failure-within-horizon 관점의 검증에 적합하다.",
            "IMS/FEMTO는 bearing run-to-failure feature 기반 조기 이상 탐지 검증에 적합하다.",
        ],
        0.9,
        4.9,
        10.6,
        1.1,
        12,
    )


def slide_product(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "Product Implementation", 20)
    add_picture(slide, OUTPUTS / "maintiq_predict_screenshot.png", 0.65, 1.3, 6.8, 3.9)
    rows = [
        ["구분", "역할", "배포 방식"],
        ["빠른 점검 모드", "작은 설치본, 경량 운영 점수", "Lite installer"],
        ["정밀 분석 모드", "XGBoost/SHAP 기반 분석", "Full installer"],
        ["Admin Console", "논문 검증, benchmark, 실증 리포트", "Local Streamlit"],
    ]
    mini_table(slide, rows, 7.7, 1.55, [2.0, 3.05, 1.65], 0.5, 8.2)
    textbox(slide, "사용자 앱은 연구 용어를 숨기고, Admin 콘솔은 검증과 논문 근거를 관리한다.", 7.8, 4.2, 4.5, 0.45, 12, C.INK, True)
    textbox(slide, "두 모드는 계산 방식이 달라 결과가 다를 수 있으며, 논문 검증은 정밀 분석/Admin 기준으로 해석한다.", 7.8, 5.05, 4.5, 0.45, 10, C.MUTED)


def slide_claim_boundary(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "Claim Boundary", 21)
    rows = [
        ["구분", "가능한 주장", "금지 주장"],
        ["구현", "설치형 제품 MVP 구현", "상용 SaaS 완제품"],
        ["검증", "공개 benchmark 기준 재현 가능 검증", "회사 로그 기반 검증이 끝났다는 주장"],
        ["운영", "승인형 작업지시 workflow", "자동 정비 명령 실행"],
        ["현장", "field validation template 준비", "현장 운영망에 배포됐다는 주장"],
        ["비용", "official/simulated cost metric 분석", "원화 비용 절감이 입증됐다는 주장"],
    ]
    mini_table(slide, rows, 0.55, 1.35, [1.65, 4.8, 4.8], 0.49, 8.2)
    claim_ready = ev.field.get("field_claim_ready", False)
    status = "현장 실증 주장 가능" if claim_ready else "현장 실증 전: 회사 로그 필요"
    color = C.GREEN2 if claim_ready else C.AMBER
    metric_card(slide, "Field validation 상태", status, f"claim_status: {ev.field.get('claim_status', '-')}", 0.85, 5.65, 6.1, accent=color)
    textbox(slide, "이 장은 발표 질의응답에서 과장 주장을 피하기 위한 방어선이다.", 7.25, 5.9, 4.4, 0.3, 11, C.MUTED)


def slide_conclusion(prs: Presentation, ev: Evidence) -> None:
    slide = blank(prs)
    title(slide, "결론", 22)
    xgb = ev.metrics["models"]["xgboost"]
    tuned = ev.threshold["selected_metrics"]
    best = next(row for row in ev.scania["metrics"] if row["strategy_id"] == "xgboost_cost_optimized")
    metric_card(slide, "모델 검증", f"PR-AUC {xgb['pr_auc']:.4f}", f"threshold tuning F1 {tuned['f1_score']:.4f}", 0.85, 1.45, 3.45, 1.2)
    metric_card(slide, "운영 통합", "예측→SPC→리포트", "승인형 작업지시까지 연결", 4.65, 1.45, 3.45, 1.2, accent=C.BLUE)
    metric_card(slide, "공개 산업 검증", pct(best["cost_improvement_vs_rule"], 2), "SCANIA official cost metric", 8.45, 1.45, 3.45, 1.2, accent=C.GREEN2)
    rows = [
        ["정리", "발표에서 말할 결론"],
        ["무엇을 만들었나", "CSV 예측, 위험 모니터링, AI 리포트, 작업지시를 통합한 데스크톱 MVP"],
        ["무엇을 검증했나", "AI4I 모델 성능, threshold trade-off, SPC 정책, SCANIA official cost metric"],
        ["무엇을 주장하지 않나", "실제 공장 배포, 실제 회사 비용 절감, 자동 정비 명령 실행"],
        ["향후 연구", "회사 labeled sensor CSV, 정비 이력, downtime/cost 로그 기반 field validation"],
    ]
    mini_table(slide, rows, 0.85, 3.15, [2.55, 8.7], 0.46, 8.5)
    shape_box(slide, 0.88, 5.95, 11.0, 0.62, C.PALE2, C.LINE)
    textbox(slide, "최종 메시지", 1.12, 6.12, 1.7, 0.22, 11, C.GREEN, True)
    textbox(slide, "본 연구의 novelty는 예측 결과를 운영 의사결정 흐름까지 연결한 재현 가능한 통합 시스템이다.", 2.6, 6.11, 8.7, 0.24, 12, C.INK, True)


def slide_references(prs: Presentation) -> None:
    slide = blank(prs)
    title(slide, "References", 23)
    refs = [
        "Jardine et al. (2006), condition-based maintenance review",
        "Carvalho et al. (2019), ML methods for predictive maintenance review",
        "Montgomery (2019), Statistical Quality Control",
        "Elkan (2001), cost-sensitive learning",
        "AI4I 2020 Predictive Maintenance Dataset; Matzka (2020)",
        "Kharazian et al. (2024), SCANIA Component X Dataset",
        "UCI MetroPT-3 Dataset; NASA C-MAPSS; NASA IMS Bearings",
        "Chen & Guestrin (2016), XGBoost",
        "Chawla et al. (2002), SMOTE",
        "Lundberg & Lee (2017), SHAP",
        "Google Gemini generateContent API; OpenAI Responses API",
        "PyInstaller documentation",
    ]
    for i, ref in enumerate(refs):
        col = 0 if i < 6 else 1
        row = i if i < 6 else i - 6
        x = 0.8 + col * 6.1
        y = 1.35 + row * 0.7
        textbox(slide, f"{i + 1}.", x, y, 0.35, 0.22, 8, C.GREEN, True)
        textbox(slide, ref, x + 0.42, y, 5.2, 0.33, 8.5, C.INK)
    textbox(slide, "상세 APA 형식 참고문헌은 outputs/research_references.md에 별도 정리했다.", 0.85, 6.35, 10.8, 0.25, 9, C.MUTED)


def build_ppt(ev: Evidence) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide_cover(prs)
    slide_agenda(prs)
    slide_motivation(prs)
    slide_contributions(prs)
    slide_review(
        prs,
        5,
        1,
        "Predictive Maintenance / CBM",
        [
            ["문헌", "핵심 내용", "본 연구 반영"],
            ["Jardine et al. (2006)", "진단, 예후, 정비 의사결정의 CBM 절차", "예측 결과를 작업지시 흐름으로 연결"],
            ["Carvalho et al. (2019)", "ML 기반 예지보전 방법론 review", "모델 비교와 운영 workflow 결합"],
            ["Montgomery (2019)", "SPC 관리도와 이상 탐지 기준", "ML 위험 확률의 Predictive SPC 구성"],
        ],
        "예지보전은 예측 모델 자체보다 정비 의사결정으로 이어지는 구조가 중요하다.",
    )
    slide_review(
        prs,
        6,
        2,
        "ML 기반 예지보전",
        [
            ["문헌", "핵심 내용", "본 연구 반영"],
            ["AI4I 2020 / Matzka", "예지보전 공개 데이터와 XAI 적용 가능성", "기준 데이터와 실험 재현성 확보"],
            ["Chen & Guestrin (2016)", "XGBoost의 비선형 feature interaction 처리", "주요 고장 예측 모델로 사용"],
            ["Chawla et al. (2002)", "SMOTE 기반 minority class 보강", "불균형 처리 비교군으로 검증"],
        ],
        "불균형 고장 데이터에서는 모델 성능과 threshold 정책을 함께 평가해야 한다.",
    )
    slide_review(
        prs,
        7,
        3,
        "SPC, SHAP, Cost-sensitive Learning",
        [
            ["문헌", "핵심 내용", "본 연구 반영"],
            ["Lundberg & Lee (2017)", "SHAP 기반 feature contribution 설명", "위험 요인 설명과 리포트 근거"],
            ["Elkan (2001)", "비용 민감 학습의 기초", "false alarm/missed failure cost simulation"],
            ["SCANIA Component X", "real-world multivariate time-series와 official cost", "공개 산업 데이터 cost metric 검증"],
        ],
        "설명 가능성과 비용 민감 평가는 예측 결과의 운영 적용성을 높이는 핵심 보완 요소다.",
    )
    slide_system(prs)
    slide_data(prs)
    slide_model(prs, ev)
    slide_spc(prs, ev)
    slide_workflow(prs, ev)
    slide_experiment_design(prs)
    slide_ai4i(prs, ev)
    slide_strategy(prs, ev)
    slide_spc_compare(prs, ev)
    slide_cost(prs, ev)
    slide_scania(prs, ev)
    slide_public_bench(prs, ev)
    slide_product(prs)
    slide_claim_boundary(prs, ev)
    slide_conclusion(prs, ev)
    slide_references(prs)
    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def get_font(size: int, bold: bool = False):
    font_path = Path(r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf")
    if font_path.exists():
        return ImageFont.truetype(str(font_path), size)
    return ImageFont.load_default()


def draw_wrapped(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, font, fill: str, width: int, line_gap: int = 5) -> int:
    x, y = xy
    max_chars = max(10, int(width / max(font.size * 0.55, 1)))
    for line in textwrap.wrap(text, width=max_chars):
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap
    return y


def preview_from_pptx() -> None:
    prs = Presentation(PPTX_PATH)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    for old in PREVIEW_DIR.glob("slide_*.png"):
        old.unlink()
    thumbs = []
    px_w, px_h = 1600, 900
    slide_w = float(prs.slide_width)
    slide_h = float(prs.slide_height)

    def to_px_rect(shape):
        return (
            int(float(shape.left) / slide_w * px_w),
            int(float(shape.top) / slide_h * px_h),
            int(float(shape.width) / slide_w * px_w),
            int(float(shape.height) / slide_h * px_h),
        )

    def fill_hex(shape, fallback: str = C.WHITE) -> str:
        try:
            fc = shape.fill.fore_color
            if fc.type is not None and fc.rgb is not None:
                return f"#{fc.rgb}"
        except Exception:
            pass
        return "#" + fallback

    def line_hex(shape, fallback: str = C.LINE) -> str:
        try:
            lc = shape.line.color
            if lc.type is not None and lc.rgb is not None:
                return f"#{lc.rgb}"
        except Exception:
            pass
        return "#" + fallback

    def draw_table(draw: ImageDraw.ImageDraw, shape) -> None:
        x, y, w, h = to_px_rect(shape)
        table_obj = shape.table
        rows = len(table_obj.rows)
        cols = len(table_obj.columns)
        if rows == 0 or cols == 0:
            return
        col_w = max(1, w // cols)
        row_h = max(1, h // rows)
        for r in range(rows):
            for c in range(cols):
                cx = x + c * col_w
                cy = y + r * row_h
                fill = "#" + (C.GREEN if r == 0 else (C.PALE2 if r % 2 else C.WHITE))
                draw.rectangle([cx, cy, cx + col_w, cy + row_h], fill=fill, outline="#" + C.LINE)
                text = table_obj.cell(r, c).text.strip().replace("\n", " ")
                if text:
                    color = "white" if r == 0 else "#" + C.INK
                    draw_wrapped(draw, (cx + 8, cy + 6), text[:85], get_font(12 if r else 13, r == 0), color, max(50, col_w - 16), 3)

    def draw_picture(img: Image.Image, shape) -> None:
        x, y, w, h = to_px_rect(shape)
        try:
            pic = Image.open(BytesIO(shape.image.blob)).convert("RGB")
            pic.thumbnail((max(1, w), max(1, h)))
            canvas = Image.new("RGB", (max(1, w), max(1, h)), "white")
            canvas.paste(pic, ((w - pic.width) // 2, (h - pic.height) // 2))
            img.paste(canvas, (x, y))
        except Exception:
            ImageDraw.Draw(img).rectangle([x, y, x + w, y + h], fill="#" + C.LIGHT, outline="#" + C.LINE)

    def draw_text_shape(draw: ImageDraw.ImageDraw, shape) -> None:
        text = shape.text.strip() if hasattr(shape, "text") else ""
        if not text:
            return
        x, y, w, h = to_px_rect(shape)
        try:
            size = int(shape.text_frame.paragraphs[0].runs[0].font.size.pt)
        except Exception:
            size = 15
        size = max(9, min(38, int(size * 1.55)))
        try:
            color = f"#{shape.text_frame.paragraphs[0].runs[0].font.color.rgb}"
        except Exception:
            color = "#" + C.INK
        draw_wrapped(draw, (x + 4, y + 2), text.replace("\n", " / "), get_font(size, False), color, max(20, w - 8), 4)

    for idx, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (px_w, px_h), "white")
        draw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                draw_table(draw, shape)
                continue
            if shape.shape_type == 13:
                draw_picture(img, shape)
                continue
            if getattr(shape, "has_chart", False):
                x, y, w, h = to_px_rect(shape)
                draw.rectangle([x, y, x + w, y + h], fill="white", outline="#" + C.LINE)
                draw.text((x + 18, y + 16), "Chart", font=get_font(18, True), fill="#" + C.GREEN)
                continue
            if hasattr(shape, "fill") and not getattr(shape, "has_text_frame", False):
                x, y, w, h = to_px_rect(shape)
                if w > 0 and h > 0:
                    draw.rectangle([x, y, x + w, y + h], fill=fill_hex(shape), outline=line_hex(shape))
            if getattr(shape, "has_text_frame", False):
                draw_text_shape(draw, shape)
        out = PREVIEW_DIR / f"slide_{idx:02d}.png"
        img.save(out)
        thumb = img.copy()
        thumb.thumbnail((360, 203))
        thumbs.append(thumb)

    cols = 4
    rows = math.ceil(len(thumbs) / cols)
    sheet = Image.new("RGB", (cols * 390, rows * 240), "white")
    sheet_draw = ImageDraw.Draw(sheet)
    for i, thumb in enumerate(thumbs):
        x = (i % cols) * 390 + 15
        y = (i // cols) * 240 + 15
        sheet.paste(thumb, (x, y))
        sheet_draw.text((x, y + 208), f"Slide {i + 1:02d}", font=get_font(16), fill="#555555")
    sheet.save(CONTACT_SHEET)


SLIDE_OUTLINE = [
    ("표지", "MaintiQ Predict 연구 제목과 핵심 메시지"),
    ("목차", "Motivation, Contributions, Review, Main Part, 실험 설계, 실험 결과, 결론"),
    ("Motivation", "예측 결과가 현장 의사결정으로 연결되지 않는 문제"),
    ("Contributions", "예측, 모니터링, 운영, 공개 benchmark 검증 기여"),
    ("Review 1: Predictive Maintenance / CBM", "CBM과 예지보전 선행연구 한계 및 반영점"),
    ("Review 2: ML 기반 예지보전", "AI4I, XGBoost, SMOTE 기반 모델링 문헌"),
    ("Review 3: SPC, SHAP, Cost-sensitive Learning", "설명 가능성과 비용 민감 평가 문헌"),
    ("Main Part: Proposed System", "CSV 입력부터 작업지시까지 전체 시스템 흐름"),
    ("Main Part: Data & Preprocessing", "AI4I 전처리와 회사 CSV 대응"),
    ("Main Part: Model and Threshold Policy", "XGBoost, threshold tuning, 주요 평가지표"),
    ("Main Part: Predictive SPC", "ML 위험 확률 기반 SPC 모니터링"),
    ("Main Part: GenAI Report and Work-order Workflow", "Gemini 리포트 실제 생성 결과와 승인형 작업지시"),
    ("실험 설계", "비교군, 데이터셋, 평가지표"),
    ("실험 결과: AI4I Baseline", "Logistic Regression과 XGBoost 성능 비교"),
    ("실험 결과: SMOTE / Threshold 비교", "불균형 처리와 threshold 정책 비교"),
    ("실험 결과: SPC-only vs ML+SPC", "경보 정책 비교"),
    ("실험 결과: Operational Cost Simulation", "false alarm/missed failure 기반 상대 cost simulation"),
    ("실험 결과: SCANIA Official Cost Metric", "SCANIA official cost matrix 검증"),
    ("실험 결과: Public Benchmark Extension", "MetroPT-3, C-MAPSS, IMS/FEMTO 확장"),
    ("Product Implementation", "Full/Lite 설치본과 Admin 콘솔"),
    ("Claim Boundary", "가능한 주장과 금지 주장"),
    ("결론", "연구 결과와 향후 실제 회사 데이터 실증 방향"),
    ("References", "핵심 참고문헌 축약 목록"),
]


def write_canva_outline() -> None:
    lines = [
        "# Canva 보조용 발표 outline",
        "",
        "Canva에서 디자인 후보를 만들 때 사용할 수 있는 23장 구성이다.",
        "스타일: 흰 배경, 진녹색 제목, 얇은 초록 제목선, 학술 발표형, 표는 초록 header와 연녹색 alternate row.",
        "",
    ]
    for idx, (title_text, description) in enumerate(SLIDE_OUTLINE, start=1):
        lines.append(f"## {idx}. {title_text}")
        lines.append(f"- {description}")
        lines.append("- 과도한 장식보다 표, 그래프, 핵심 문장 중심으로 구성")
        lines.append("")
    CANVA_OUTLINE_PATH.write_text("\n".join(lines), encoding="utf-8")


def write_notes(ev: Evidence) -> None:
    xgb = ev.metrics["models"]["xgboost"]
    tuned = ev.threshold["selected_metrics"]
    scania_best = next(row for row in ev.scania["metrics"] if row["strategy_id"] == "xgboost_cost_optimized")
    ai = ai_report_summary(ev)
    NOTES_PATH.write_text(
        f"""# MaintiQ Predict 최종발표 발표자 노트

## 전체 발표 메시지
- 본 연구는 단일 고장 예측 모델이 아니라 `예측 → SPC → GenAI 리포트 → 승인형 작업지시 → 공개 benchmark 검증`을 연결한 통합 시스템이다.
- AI4I 기준 XGBoost PR-AUC는 {xgb['pr_auc']:.4f}, threshold tuning 후 F1은 {tuned['f1_score']:.4f}이다.
- Gemini 리포트는 `{ai['mode']}`로 생성되었고, UDI {ai['udi']} row에 대해 예측 확률 {ai['probability']}, 기준 {ai['threshold']}, 상태 {ai['risk_status']}를 요약했다.
- SCANIA 공개 benchmark에서는 official cost metric 기준 rule baseline 대비 {scania_best['cost_improvement_vs_rule'] * 100:.2f}% 개선 가능성을 확인했다.
- 실제 회사 비용 절감이나 실제 공장 배포 완료가 아니라, 공개 데이터와 로컬 제품형 MVP 기반 검증으로 한정한다.

## 슬라이드별 말할 포인트
1. 표지: 시스템 이름과 연구 범위를 먼저 말한다. 실제 배포가 아닌 제품형 MVP임을 짧게 못박는다.
2. 목차: 7개 흐름으로 발표가 진행된다고 안내한다.
3. Motivation: 예측 결과가 작업지시로 이어지지 않으면 현장 가치가 낮다는 문제를 강조한다.
4. Contributions: 모델, SPC, GenAI, 작업지시, 공개 benchmark의 통합이 기여임을 설명한다.
5~7. Review: 선행연구를 나열하지 말고 본 연구에 어떻게 반영했는지 중심으로 설명한다.
8~12. Main Part: CSV 입력부터 리포트와 작업지시까지 실제 시스템 흐름을 연결해서 설명한다. 특히 12번 슬라이드는 실제 Gemini 리포트 생성 evidence를 보여준다.
13. 실험 설계: 비교군과 지표를 먼저 보여줘 결과 해석 기준을 세운다.
14~19. 실험 결과: 각 결과는 한 문장 해석만 강조하고 세부표는 질문 때 설명한다.
20. Product Implementation: Full/Lite와 Admin 역할을 구분한다.
21. Claim Boundary: 가능한 주장과 금지 주장을 명확히 구분해 과장 우려를 차단한다.
22. 결론: 시스템 통합성과 향후 실제 회사 데이터 실증 계획을 강조한다.
23. References: 핵심 문헌만 보여주고 상세 참고문헌은 별도 문서에 있다고 말한다.

## 발표 중 금지할 표현
- 실제 PLC/SCADA 운영망 배포 완료
- 회사 로그 기반 성능 검증이 끝났다는 표현
- 원화 비용 절감이 입증됐다는 표현
- 자동 정비를 시스템이 직접 수행했다는 표현
""",
        encoding="utf-8",
    )


def validate_outputs() -> None:
    prs = Presentation(PPTX_PATH)
    titles = []
    for slide in prs.slides:
        slide_texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        titles.append(slide_texts[0].splitlines()[0] if slide_texts else "")
    required = [
        "Review 1: Predictive Maintenance / CBM",
        "Review 2: ML 기반 예지보전",
        "Review 3: SPC, SHAP, Cost-sensitive Learning",
        "Main Part: Proposed System",
        "실험 설계",
        "실험 결과: AI4I Baseline",
        "결론",
        "References",
    ]
    missing = [item for item in required if item not in titles]
    if len(prs.slides) != 23:
        raise AssertionError(f"Expected 23 slides, found {len(prs.slides)}")
    if missing:
        raise AssertionError(f"Missing required slide titles: {missing}")
    forbidden = [
        "실제 PLC/SCADA 배포 완료",
        "실제 회사 데이터 검증 완료",
        "실제 비용 절감 실증 완료",
        "자동 정비 명령 실행",
    ]
    combined = "\n".join(
        path.read_text(encoding="utf-8", errors="replace")
        for path in [
            OUTLINE_PATH,
            TOC_PATH,
            REF_PATH,
            NOTES_PATH,
            CANVA_OUTLINE_PATH,
            GENAI_EVIDENCE_PATH,
            PROFESSOR_BRIEF_PATH,
        ]
        if path.exists()
    )
    ppt_text = "\n".join(titles)
    for phrase in forbidden:
        if phrase in combined or phrase in ppt_text:
            raise AssertionError(f"Forbidden claim found: {phrase}")
    full_ppt_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )
    for needle in ["gemini_generate_content:gemini-2.5-flash", "0.993616", "0.87", "High Risk"]:
        if needle not in full_ppt_text:
            raise AssertionError(f"GenAI evidence missing from PPT text: {needle}")
    for path in [OUTLINE_PATH, TOC_PATH, REF_PATH, NOTES_PATH, CANVA_OUTLINE_PATH, GENAI_EVIDENCE_PATH, PROFESSOR_BRIEF_PATH]:
        text = path.read_text(encoding="utf-8", errors="replace")
        if "\ufffd" in text or "??" in text:
            raise AssertionError(f"Encoding issue suspected in {path}")


def main() -> None:
    ev = load_evidence()
    OUTPUTS.mkdir(parents=True, exist_ok=True)
    write_references()
    write_thesis_docs(ev)
    write_toc()
    write_genai_evidence_docs(ev)
    build_ppt(ev)
    preview_from_pptx()
    write_canva_outline()
    write_notes(ev)
    validate_outputs()
    print(f"Thesis outline: {OUTLINE_PATH}")
    print(f"Thesis TOC: {TOC_PATH}")
    print(f"References: {REF_PATH}")
    print(f"PPTX: {PPTX_PATH}")
    print(f"Preview: {CONTACT_SHEET}")
    print(f"Canva outline: {CANVA_OUTLINE_PATH}")
    print(f"Notes: {NOTES_PATH}")
    print(f"GenAI evidence: {GENAI_EVIDENCE_PATH}")
    print(f"Professor brief: {PROFESSOR_BRIEF_PATH}")


if __name__ == "__main__":
    main()
