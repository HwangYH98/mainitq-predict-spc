"""Create the final presentation deck for MaintiQ Predict.

The deck is built from local, already-generated validation artifacts.  It keeps
the claims defensible: public benchmark and simulation evidence is separated
from real field proof, and no slide claims real factory deployment or real cost
reduction without company logs.
"""

from __future__ import annotations

import csv
import json
import math
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = ROOT / "outputs"
PREVIEW_DIR = OUTPUTS / "final_presentation_preview"
PPTX_PATH = OUTPUTS / "final_presentation_maintiq_predict.pptx"
CONTACT_SHEET_PATH = PREVIEW_DIR / "contact_sheet.png"
SPEAKER_NOTES_PATH = OUTPUTS / "final_presentation_speaker_notes.md"

SLIDE_W = 13.333
SLIDE_H = 7.5


class C:
    NAVY = "12233A"
    INK = "0B1220"
    MUTED = "5E6A7D"
    LINE = "D8E1EC"
    BG = "F4F7FB"
    WHITE = "FFFFFF"
    BLUE = "246B9F"
    CYAN = "2CA6A4"
    GREEN = "2E9D69"
    AMBER = "E59B2D"
    RED = "D84B4B"
    PALE_BLUE = "E9F3FB"
    PALE_GREEN = "EAF7F1"
    PALE_AMBER = "FFF5E4"


FONT = "Malgun Gothic"


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def load_json(path: str) -> dict:
    return json.loads((ROOT / path).read_text(encoding="utf-8"))


def load_csv(path: str) -> list[dict[str, str]]:
    with (ROOT / path).open("r", encoding="utf-8-sig", newline="") as handle:
        return list(csv.DictReader(handle))


def pct(value: float | str, digits: int = 1) -> str:
    return f"{float(value) * 100:.{digits}f}%"


def num(value: float | str, digits: int = 3) -> str:
    return f"{float(value):.{digits}f}"


def text_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 18,
    color: str = C.INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    fill: str | None = None,
    line: str | None = None,
    radius: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill or C.WHITE)
    if fill is None:
        shape.fill.transparency = 100000
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.transparency = 100000
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_title(slide, title: str, subtitle: str | None = None):
    text_box(slide, title, 0.55, 0.35, 8.8, 0.48, 22, C.INK, True)
    if subtitle:
        text_box(slide, subtitle, 0.58, 0.87, 10.2, 0.35, 9, C.MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(C.LINE)
    line.line.transparency = 100000


def add_footer(slide, idx: int, note: str = "MaintiQ Predict final presentation"):
    text_box(slide, note, 0.55, 7.12, 7.5, 0.22, 7, C.MUTED)
    text_box(slide, f"{idx:02d}", 12.25, 7.1, 0.45, 0.22, 7, C.MUTED, align=PP_ALIGN.RIGHT)


def add_metric(slide, label: str, value: str, caption: str, x: float, y: float, w: float, color: str = C.BLUE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(C.WHITE)
    box.line.color.rgb = rgb(C.LINE)
    text_box(slide, label, x + 0.15, y + 0.12, w - 0.3, 0.22, 8, C.MUTED, True)
    text_box(slide, value, x + 0.15, y + 0.36, w - 0.3, 0.34, 18, color, True)
    text_box(slide, caption, x + 0.15, y + 0.73, w - 0.3, 0.22, 7, C.MUTED)


def add_bullet_list(slide, items: Iterable[str], x: float, y: float, w: float, h: float, font_size: int = 12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb(C.INK)
        p.space_after = Pt(7)
    return box


def add_image(slide, path: Path, x: float, y: float, w: float, h: float):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        text_box(slide, f"이미지 없음: {path.name}", x, y, w, h, 10, C.RED, fill=C.WHITE, line=C.LINE, radius=True)


def add_chart(slide, title: str, categories: list[str], series: list[tuple[str, list[float]]], x: float, y: float, w: float, h: float):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        chart_data,
    )
    chart = graphic.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = title
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = 1.0
    chart.value_axis.minimum_scale = 0.0
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT
    return chart


@dataclass
class Evidence:
    metrics: dict
    threshold: dict
    spc: dict
    calibration: dict
    model_rows: list[dict[str, str]]
    spc_rows: list[dict[str, str]]
    op_rows: list[dict[str, str]]
    scania: dict
    field: dict
    public_rows: list[dict[str, str]]


def load_evidence() -> Evidence:
    return Evidence(
        metrics=load_json("outputs/metrics.json"),
        threshold=load_json("outputs/threshold_summary.json"),
        spc=load_json("outputs/spc_summary.json"),
        calibration=load_json("outputs/probability_calibration_metrics.json"),
        model_rows=load_csv("outputs/model_strategy_comparison.csv"),
        spc_rows=load_csv("outputs/spc_vs_ml_comparison.csv"),
        op_rows=load_csv("outputs/operational_value_simulation.csv"),
        scania=load_json("outputs/scania_official_cost_metrics.json"),
        field=load_json("outputs/field_validation_report.json"),
        public_rows=load_csv("outputs/public_industrial_validation_metrics.csv"),
    )


def slide_cover(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(C.BG)
    # editorial side rail
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.05), Inches(SLIDE_H))
    rail.fill.solid()
    rail.fill.fore_color.rgb = rgb(C.NAVY)
    rail.line.transparency = 100000
    text_box(slide, "MaintiQ\nPredict", 0.55, 0.78, 3.05, 1.35, 34, C.WHITE, True)
    text_box(slide, "AI 예지보전 운영 대시보드\n최종발표", 0.62, 2.32, 2.9, 0.65, 14, "CFE3F5", True)
    text_box(slide, "CSV 예측 · SPC 위험 분석 · GenAI 리포트 · 승인형 작업지시", 0.62, 6.48, 3.05, 0.35, 9, "D8E7F5")
    text_box(slide, "AI 예지보전 시스템을\n실행 가능한 제품형 MVP로 구현하고\n공개 산업 데이터로 검증했습니다.", 4.6, 0.88, 7.6, 1.55, 27, C.INK, True)
    text_box(
        slide,
        "핵심 메시지: 단순 예측 모델이 아니라, 전처리·위험판정·SPC·AI 리포트·작업지시·검증 근거를 한 흐름으로 연결한 로컬 운영 시스템입니다.",
        4.65,
        2.55,
        7.5,
        0.75,
        13,
        C.MUTED,
    )
    add_image(slide, OUTPUTS / "maintiq_predict_screenshot.png", 4.65, 3.55, 7.8, 2.95)
    add_footer(slide, 1, "Final presentation deck")


def slide_problem(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "문제 정의: 예측은 있어도 운영 의사결정까지 이어지기 어렵다")
    steps = [
        ("사후보전", "고장이 난 뒤 대응하면 downtime과 정비 비용이 커짐", C.RED),
        ("단일 임계값/SPC", "한 센서나 고정 기준만 보면 복합 센서 조합을 놓칠 수 있음", C.AMBER),
        ("설명·조치 단절", "예측 결과가 작업지시, 승인 기록, 재학습 후보로 남지 않음", C.BLUE),
    ]
    for i, (head, body, color) in enumerate(steps):
        x = 0.8 + i * 4.05
        text_box(slide, head, x, 1.65, 3.2, 0.36, 22, color, True)
        text_box(slide, body, x, 2.15, 3.25, 1.15, 16, C.INK)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.48), Inches(3.1), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(color)
        line.line.transparency = 100000
    text_box(slide, "해결 방향", 0.82, 4.35, 2.2, 0.35, 18, C.INK, True)
    add_bullet_list(
        slide,
        [
            "CSV 입력부터 전처리, 위험 확률, SPC 그래프, AI 리포트, 작업지시 승인까지 한 화면 흐름으로 연결",
            "실제 현장 배포 완료가 아니라, 재현 가능한 제품형 MVP와 공개 benchmark 검증을 명확히 분리",
            "비용 절감은 실제 회사 로그 없이는 주장하지 않고, simulation/official cost metric으로만 표현",
        ],
        0.85,
        4.85,
        11.6,
        1.25,
        13,
    )
    add_footer(slide, 2)


def slide_architecture(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "시스템 구조: 데이터에서 작업지시까지 이어지는 폐쇄 루프")
    nodes = [
        ("센서 CSV", "AI4I형 / 회사형 CSV\n컬럼 매핑·단위 변환", C.PALE_BLUE, C.BLUE),
        ("전처리·품질 진단", "결측·숫자 오류·분포 이탈\n품질 리포트 생성", C.WHITE, C.BLUE),
        ("위험 예측", "XGBoost/threshold\ncalibrated probability", C.PALE_GREEN, C.GREEN),
        ("SPC·모니터링", "관리한계·위험 추세\n고위험 설비 우선순위", C.WHITE, C.CYAN),
        ("AI 리포트", "Gemini/OpenAI\n관리자 참고 리포트", C.PALE_AMBER, C.AMBER),
        ("작업지시", "초안 생성\n승인/검토/반려 기록", C.WHITE, C.RED),
    ]
    y = 1.65
    for i, (head, body, fill, color) in enumerate(nodes):
        x = 0.55 + i * 2.08
        slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.82), Inches(1.55)).fill.solid()
        shape = slide.shapes[-1]
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(C.LINE)
        text_box(slide, head, x + 0.12, y + 0.15, 1.58, 0.28, 12, color, True)
        text_box(slide, body, x + 0.12, y + 0.55, 1.58, 0.66, 8, C.MUTED)
        if i < len(nodes) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.82), Inches(y + 0.78), Inches(x + 2.06), Inches(y + 0.78))
            conn.line.color.rgb = rgb(C.MUTED)
            conn.line.width = Pt(1.5)
    text_box(slide, "운영 이력", 0.8, 4.25, 2.3, 0.38, 18, C.INK, True)
    add_bullet_list(
        slide,
        [
            "운영 이벤트, 작업지시 초안, 작업자 결정은 SQLite와 CSV export로 추적",
            "needs_review 결정은 재학습 후보로 남겨 모델 개선 흐름과 연결",
            "Admin 콘솔은 모델 비교, SCANIA/MetroPT 공개 benchmark, field validation template을 별도로 관리",
        ],
        0.85,
        4.82,
        6.0,
        1.35,
        12,
    )
    text_box(slide, "배포 형태", 7.45, 4.25, 2.3, 0.38, 18, C.INK, True)
    add_bullet_list(
        slide,
        [
            "Full: 정밀 분석 모드, XGBoost/SHAP 경로 유지",
            "Lite: 빠른 점검 모드, 32.8MB 경량 설치본",
            "설치파일은 GitHub Release 첨부, repo에는 코드/샘플/검증 스크립트 중심",
        ],
        7.5,
        4.82,
        5.2,
        1.35,
        12,
    )
    add_footer(slide, 3)


def slide_product(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "제품 화면: 브라우저 없이 실행되는 Windows 데스크톱 앱")
    add_image(slide, OUTPUTS / "maintiq_predict_screenshot.png", 0.65, 1.45, 8.0, 4.55)
    add_metric(slide, "Full", "정밀 분석 모드", "XGBoost/SHAP 기반", 9.05, 1.55, 3.35, C.BLUE)
    add_metric(slide, "Lite", "빠른 점검 모드", "경량 운영 점수", 9.05, 2.85, 3.35, C.GREEN)
    add_metric(slide, "설치본", "32.8MB / 206.4MB", "Lite / Full 기준", 9.05, 4.15, 3.35, C.AMBER)
    text_box(
        slide,
        "두 모드는 계산 방식이 다르므로 결과가 다를 수 있습니다. 최종 검증과 논문 근거는 Full/Admin 경로를 기준으로 확인합니다.",
        9.08,
        5.65,
        3.25,
        0.75,
        10,
        C.MUTED,
    )
    add_footer(slide, 4)


def slide_data_pipeline(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "데이터 처리: 학습 누수 제거와 회사 CSV 대응을 분리")
    add_metric(slide, "학습 데이터", "AI4I 2020", "10,000 rows", 0.7, 1.55, 3.0, C.BLUE)
    add_metric(slide, "테스트 분할", "2,000 rows", "stratified split", 3.95, 1.55, 3.0, C.CYAN)
    add_metric(slide, "특징 수", "8 features", "Type one-hot 포함", 7.2, 1.55, 3.0, C.GREEN)
    add_metric(slide, "확률 보정", "Isotonic", "Brier 0.0124", 10.45, 1.55, 2.2, C.AMBER)
    text_box(slide, "전처리 원칙", 0.78, 3.2, 2.4, 0.32, 16, C.INK, True)
    add_bullet_list(
        slide,
        [
            "목표값: Machine failure",
            "ID성 컬럼 제거: UDI, Product ID",
            "누수 컬럼 제거: TWF, HDF, PWF, OSF, RNF",
            "Type은 one-hot encoding 처리",
        ],
        0.8,
        3.68,
        5.0,
        1.55,
        12,
    )
    text_box(slide, "회사 CSV 입력 보강", 6.6, 3.2, 2.9, 0.32, 16, C.INK, True)
    add_bullet_list(
        slide,
        [
            "컬럼 alias 자동 매핑: rpm, 온도, torque 등",
            "Celsius/Kelvin, percent/ratio 단위 변환",
            "결측률, 숫자 변환 실패, 허용 범위 밖 값 진단",
            "결과 CSV에는 engine_profile, score_method, interpretation_note 기록",
        ],
        6.62,
        3.68,
        5.75,
        1.55,
        12,
    )
    add_footer(slide, 5)


def slide_model_metrics(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "모델 결과: XGBoost가 baseline보다 PR-AUC와 F1에서 우세")
    models = ev.metrics["models"]
    add_chart(
        slide,
        "Baseline comparison",
        ["Precision", "Recall", "F1", "ROC-AUC", "PR-AUC"],
        [
            (
                "Logistic Regression",
                [
                    models["logistic_regression"]["precision"],
                    models["logistic_regression"]["recall"],
                    models["logistic_regression"]["f1_score"],
                    models["logistic_regression"]["roc_auc"],
                    models["logistic_regression"]["pr_auc"],
                ],
            ),
            (
                "XGBoost",
                [
                    models["xgboost"]["precision"],
                    models["xgboost"]["recall"],
                    models["xgboost"]["f1_score"],
                    models["xgboost"]["roc_auc"],
                    models["xgboost"]["pr_auc"],
                ],
            ),
        ],
        0.68,
        1.55,
        7.2,
        4.3,
    )
    th = ev.threshold
    add_metric(slide, "선택 threshold", f"{th['selected_threshold']:.2f}", "F1 기준 선택", 8.35, 1.62, 3.75, C.BLUE)
    add_metric(slide, "Tuned F1", num(th["selected_metrics"]["f1_score"], 4), "precision/recall 균형", 8.35, 2.95, 3.75, C.GREEN)
    add_metric(slide, "Tuned precision", num(th["selected_metrics"]["precision"], 4), "오경보 감소", 8.35, 4.28, 3.75, C.CYAN)
    text_box(slide, "해석", 8.38, 5.68, 0.8, 0.28, 11, C.INK, True)
    text_box(
        slide,
        "기본 0.5 threshold는 recall이 높지만 오경보가 많습니다. 발표에서는 0.87 threshold를 운영 정책 기준으로 설명합니다.",
        9.05,
        5.63,
        3.2,
        0.58,
        9,
        C.MUTED,
    )
    add_footer(slide, 6)


def slide_alert_policy(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "위험판정 정책: SPC-only보다 ML threshold가 더 균형적")
    categories = [r["display_name"].replace(" torque control limit", "").replace(" selected threshold", "") for r in ev.spc_rows]
    add_chart(
        slide,
        "Alert policy comparison",
        categories,
        [
            ("Precision", [float(r["precision"]) for r in ev.spc_rows]),
            ("Recall", [float(r["recall"]) for r in ev.spc_rows]),
            ("F1", [float(r["f1_score"]) for r in ev.spc_rows]),
        ],
        0.65,
        1.45,
        7.6,
        4.45,
    )
    for i, row in enumerate(ev.spc_rows):
        add_metric(
            slide,
            row["display_name"].split(" ")[0],
            row["alert_count"],
            f"alerts / F1 {float(row['f1_score']):.3f}",
            8.65,
            1.45 + i * 1.42,
            3.35,
            [C.AMBER, C.BLUE, C.GREEN][i],
        )
    text_box(
        slide,
        "논문용 포인트: 단일 센서 SPC rule은 precision은 높지만 recall이 낮고, ML+SPC는 놓치는 고장을 줄이는 방향의 운영 정책으로 설명할 수 있습니다.",
        8.7,
        5.88,
        3.3,
        0.62,
        9,
        C.MUTED,
    )
    add_footer(slide, 7)


def slide_operational_value(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "운영 가치 평가는 실제 비용이 아니라 normalized cost simulation")
    conservative = [r for r in ev.op_rows if r["scenario_id"] == "conservative"]
    categories = [r["display_name"].replace(" threshold", "").replace(" baseline", "") for r in conservative]
    values = [float(r["normalized_operating_cost"]) for r in conservative]
    add_chart(slide, "Conservative scenario: normalized operating cost", categories, [("Cost", values)], 0.72, 1.5, 7.25, 4.45)
    best = min(conservative, key=lambda r: float(r["normalized_operating_cost"]))
    add_metric(slide, "최저 simulation cost", num(best["normalized_operating_cost"], 3), best["display_name"], 8.4, 1.75, 3.65, C.GREEN)
    add_metric(slide, "No-alert 기준", "1.000", "reference baseline", 8.4, 3.1, 3.65, C.MUTED)
    text_box(slide, "주의", 8.45, 4.65, 1.0, 0.3, 14, C.RED, True)
    text_box(
        slide,
        "이 값은 false alarm, missed failure, planned action에 상대 비용을 둔 simulation입니다. 실제 원화 절감률이나 현장 downtime 감소 실증으로 표현하지 않습니다.",
        8.45,
        5.05,
        3.75,
        0.9,
        10,
        C.MUTED,
    )
    add_footer(slide, 8)


def slide_scaina(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "공개 산업 데이터 검증: SCANIA official cost metric")
    best = next(r for r in ev.scania["metrics"] if r["strategy_id"] == "xgboost_cost_optimized")
    add_metric(slide, "데이터셋", "SCANIA Component X", "실제 fleet 기반 공개 benchmark", 0.7, 1.45, 3.7, C.BLUE)
    add_metric(slide, "공식 class", "0~4", "48~0 time-unit window", 4.65, 1.45, 2.4, C.CYAN)
    add_metric(slide, "Rule 대비", pct(best["cost_improvement_vs_rule"], 2), "official cost metric 개선", 7.3, 1.45, 2.35, C.GREEN)
    add_metric(slide, "No-alert 대비", pct(best["cost_improvement_vs_no_alert"], 2), "official cost metric 개선", 9.9, 1.45, 2.45, C.GREEN)
    add_image(slide, OUTPUTS / "scania_official_cost_comparison.png", 0.75, 3.05, 5.7, 2.9)
    add_image(slide, OUTPUTS / "scania_official_confusion_matrix.png", 6.85, 3.05, 5.1, 2.9)
    text_box(
        slide,
        "주장 경계: 이 결과는 SCANIA 공개 benchmark의 공식 cost matrix 기준입니다. 실제 회사의 원화 비용 절감 또는 현장 downtime 감소 실증으로 말하면 안 됩니다.",
        0.78,
        6.25,
        11.6,
        0.45,
        10,
        C.MUTED,
    )
    add_footer(slide, 9)


def slide_public_benchmarks(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "공개 benchmark 확장: 데이터 유형별 검증 경로를 준비")
    datasets = {}
    for row in ev.public_rows:
        datasets.setdefault(row["dataset_id"], []).append(row)
    rows = []
    for dataset_id, group in datasets.items():
        xgb = next((r for r in group if r["strategy_id"] == "xgboost_tuned_threshold"), group[-1])
        rows.append([dataset_id.upper(), xgb["dataset_name"], xgb["source_mode"], f"F1 {float(xgb['f1_score']):.3f}", f"PR-AUC {float(xgb['pr_auc']):.3f}"])
    table = slide.shapes.add_table(len(rows) + 1, 5, Inches(0.65), Inches(1.6), Inches(12.0), Inches(3.15)).table
    headers = ["Dataset", "Name", "Mode", "F1", "PR-AUC"]
    col_widths = [1.25, 3.8, 3.35, 1.4, 1.4]
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(C.NAVY)
        cell.text_frame.paragraphs[0].font.color.rgb = rgb(C.WHITE)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.name = FONT
        cell.text_frame.paragraphs[0].font.size = Pt(8)
    for r, values in enumerate(rows[:7], start=1):
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.name = FONT
            cell.text_frame.paragraphs[0].font.size = Pt(7.5)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(C.WHITE if r % 2 else "F8FAFD")
    text_box(slide, "해석", 0.72, 5.2, 1.0, 0.3, 16, C.INK, True)
    add_bullet_list(
        slide,
        [
            "MetroPT-3는 압축기 이상/고장 horizon 분석에 적합하지만 비용 로그는 없음",
            "C-MAPSS/IMS/FEMTO는 RUL·run-to-failure·early warning 검증을 보강하는 데이터",
            "실제 회사 효과 주장은 field sensor CSV + 정비 이력 + downtime/cost 로그가 들어온 뒤에만 가능",
        ],
        0.75,
        5.68,
        11.6,
        0.9,
        11,
    )
    add_footer(slide, 10)


def slide_boundaries(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "적용 범위: 구현된 것과 아직 주장하면 안 되는 것")
    left = [
        "Windows 데스크톱 제품형 MVP",
        "AI4I 기반 예측·threshold·SPC·calibration",
        "Gemini/OpenAI 세션형 리포트 생성",
        "승인형 작업지시 workflow와 감사 이력",
        "SCANIA official cost metric 공개 benchmark",
        "회사 field validation 템플릿과 export",
    ]
    right = [
        "실제 PLC/SCADA 운영망 연결 완료",
        "실제 공장 센서 실시간 운영 배포",
        "실제 회사 데이터 성능 재검증 완료",
        "실제 비용 절감률·탐지 시간 단축률 실증",
        "자동 정비 명령 실행",
        "상용 SaaS 완제품 수준",
    ]
    text_box(slide, "발표에서 말할 수 있음", 0.8, 1.55, 4.8, 0.38, 17, C.GREEN, True)
    add_bullet_list(slide, left, 0.95, 2.1, 5.35, 3.45, 13)
    text_box(slide, "외부 증거 없이는 금지", 7.0, 1.55, 4.8, 0.38, 17, C.RED, True)
    add_bullet_list(slide, right, 7.15, 2.1, 5.35, 3.45, 13)
    text_box(
        slide,
        f"현재 field validation 상태: {ev.field.get('claim_status', 'unknown')} / field_claim_ready={ev.field.get('field_claim_ready')}",
        0.82,
        6.25,
        11.2,
        0.35,
        10,
        C.MUTED,
    )
    add_footer(slide, 11)


def slide_conclusion(prs: Presentation, ev: Evidence):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "결론: 연구의 새로움은 모델 하나가 아니라 운영 흐름의 통합")
    text_box(slide, "Something New", 0.8, 1.35, 4.0, 0.55, 24, C.BLUE, True)
    text_box(
        slide,
        "예측 모델 결과를 threshold, SPC, GenAI 설명, 승인형 작업지시, 공개 benchmark 검증까지 연결한 재현 가능한 경량 운영 시스템",
        0.82,
        2.05,
        11.5,
        0.78,
        19,
        C.INK,
        True,
    )
    contributions = [
        ("1", "예측", "XGBoost + threshold tuning으로 불균형 failure 예측 성능 개선"),
        ("2", "운영", "SPC와 ML 확률을 결합해 경보 정책을 비교하고 작업지시로 연결"),
        ("3", "설명", "SHAP/위험요인/GenAI 리포트로 관리자 의사결정 보조"),
        ("4", "검증", "AI4I, SCANIA, MetroPT-3 등 공개 benchmark와 claim guardrail 구축"),
    ]
    for i, (no, head, body) in enumerate(contributions):
        x = 0.9 + (i % 2) * 5.9
        y = 3.35 + (i // 2) * 1.25
        text_box(slide, no, x, y, 0.45, 0.45, 18, C.WHITE, True, align=PP_ALIGN.CENTER, fill=C.BLUE, radius=True)
        text_box(slide, head, x + 0.65, y - 0.02, 1.3, 0.25, 13, C.INK, True)
        text_box(slide, body, x + 0.65, y + 0.33, 4.7, 0.35, 10, C.MUTED)
    text_box(
        slide,
        "다음 단계는 실제 회사 로그 확보입니다. labeled sensor CSV, 정비 이력, downtime/cost 로그가 들어와야 실제 비용 절감과 탐지 시간 단축을 주장할 수 있습니다.",
        0.85,
        6.15,
        11.6,
        0.45,
        11,
        C.MUTED,
    )
    add_footer(slide, 12)


def build_pptx(ev: Evidence) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    # Blank default decks often keep an empty title slide; remove all initial slides.
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
    for fn in [
        slide_cover,
        slide_problem,
        slide_architecture,
        slide_product,
        slide_data_pipeline,
        slide_model_metrics,
        slide_alert_policy,
        slide_operational_value,
        slide_scaina,
        slide_public_benchmarks,
        slide_boundaries,
        slide_conclusion,
    ]:
        fn(prs, ev)
    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def get_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path(r"C:\Windows\Fonts\malgunbd.ttf") if bold else Path(r"C:\Windows\Fonts\malgun.ttf"),
        Path(r"C:\Windows\Fonts\arialbd.ttf") if bold else Path(r"C:\Windows\Fonts\arial.ttf"),
    ]
    for path in candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def draw_wrapped(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, font, fill: str, width: int, line_gap: int = 6):
    x, y = xy
    avg = max(1, int(width / max(8, font.size * 0.62)))
    for para in text.split("\n"):
        for line in textwrap.wrap(para, width=avg):
            draw.text((x, y), line, font=font, fill=fill)
            y += font.size + line_gap
        y += line_gap
    return y


def preview_slide(number: int, title: str, subtitle: str, accents: list[tuple[str, str]], image: Path | None = None):
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    canvas = Image.new("RGB", (1600, 900), "#" + C.BG)
    draw = ImageDraw.Draw(canvas)
    draw.rectangle((0, 0, 1600, 900), fill="#" + C.BG)
    draw.text((70, 55), title, font=get_font(44, True), fill="#" + C.INK)
    draw_wrapped(draw, (72, 122), subtitle, get_font(22), "#" + C.MUTED, 1160)
    y = 220
    for label, value in accents[:5]:
        draw.rounded_rectangle((80, y, 610, y + 95), radius=18, fill="#FFFFFF", outline="#" + C.LINE, width=2)
        draw.text((105, y + 18), label, font=get_font(20, True), fill="#" + C.MUTED)
        draw_wrapped(draw, (105, y + 47), value, get_font(21, True), "#" + C.INK, 470, 3)
        y += 118
    if image and image.exists():
        try:
            img = Image.open(image).convert("RGB")
            img.thumbnail((840, 520))
            x = 700
            yy = 250
            draw.rounded_rectangle((x - 15, yy - 15, x + img.width + 15, yy + img.height + 15), radius=20, fill="#FFFFFF", outline="#" + C.LINE, width=2)
            canvas.paste(img, (x, yy))
        except Exception as exc:  # noqa: BLE001
            draw.text((700, 300), f"image preview failed: {exc}", font=get_font(20), fill="#D84B4B")
    draw.text((1480, 840), f"{number:02d}", font=get_font(20, True), fill="#" + C.MUTED)
    out = PREVIEW_DIR / f"slide_{number:02d}.png"
    canvas.save(out)
    return out


def build_previews(ev: Evidence) -> None:
    slides = [
        ("MaintiQ Predict", "AI 예지보전 운영 대시보드 최종발표", [("핵심", "제품형 MVP + 공개 산업 데이터 검증")], OUTPUTS / "maintiq_predict_screenshot.png"),
        ("문제 정의", "예측 결과가 실제 운영 의사결정으로 이어지지 않는 문제", [("해결", "예측·SPC·AI 리포트·작업지시 연결")], None),
        ("시스템 구조", "CSV 입력부터 승인형 작업지시까지 이어지는 폐쇄 루프", [("구성", "전처리 → 예측 → SPC → AI → 작업지시")], None),
        ("제품 화면", "브라우저 없이 실행되는 Windows 데스크톱 앱", [("Full", "정밀 분석 모드"), ("Lite", "빠른 점검 모드")], OUTPUTS / "maintiq_predict_screenshot.png"),
        ("데이터 처리", "학습 누수 제거와 회사 CSV 대응을 분리", [("데이터", "AI4I 2020 / 10,000 rows"), ("품질", "컬럼 매핑·단위 변환·결측 진단")], None),
        ("모델 결과", "XGBoost가 baseline보다 PR-AUC와 F1에서 우세", [("PR-AUC", "XGBoost 0.8014"), ("Threshold", "0.87 / F1 0.7752")], None),
        ("경보 정책", "SPC-only보다 ML threshold가 균형적", [("SPC-only", "F1 0.160"), ("ML threshold", "F1 0.775"), ("ML+SPC", "Recall 0.809")], None),
        ("운영 가치 simulation", "실제 비용이 아닌 normalized cost simulation", [("Best", "ML+SPC normalized cost 0.454"), ("주의", "실제 비용 절감 실증 아님")], OUTPUTS / "operational_value_simulation.png"),
        ("SCANIA benchmark", "공식 cost matrix 기준 공개 산업 데이터 검증", [("Rule 대비", "official cost metric 17.02% 개선"), ("경계", "실제 원화 비용 절감 아님")], OUTPUTS / "scania_official_cost_comparison.png"),
        ("공개 benchmark 확장", "MetroPT-3, C-MAPSS, IMS/FEMTO 검증 경로", [("MetroPT proxy", "XGBoost tuned F1 0.841"), ("목적", "lead-time/RUL/evidence 보강")], None),
        ("적용 범위", "말할 수 있는 것과 금지할 주장을 분리", [("가능", "제품형 MVP와 공개 benchmark"), ("금지", "실제 현장 비용 절감 실증")], None),
        ("결론", "새로움은 모델 하나가 아니라 운영 흐름의 통합", [("Contribution", "ML+SPC+GenAI+작업지시+검증")], None),
    ]
    paths = [preview_slide(i + 1, *spec) for i, spec in enumerate(slides)]
    # Contact sheet
    thumbs = []
    for p in paths:
        img = Image.open(p).convert("RGB")
        img.thumbnail((360, 203))
        thumbs.append(img.copy())
    sheet = Image.new("RGB", (4 * 390, 3 * 245), "#FFFFFF")
    draw = ImageDraw.Draw(sheet)
    for idx, img in enumerate(thumbs):
        x = (idx % 4) * 390 + 15
        y = (idx // 4) * 245 + 15
        sheet.paste(img, (x, y))
        draw.text((x, y + 208), f"Slide {idx + 1:02d}", font=get_font(16, True), fill="#5E6A7D")
    sheet.save(CONTACT_SHEET_PATH)


def write_speaker_notes(ev: Evidence) -> None:
    th = ev.threshold
    best_scania = next(r for r in ev.scania["metrics"] if r["strategy_id"] == "xgboost_cost_optimized")
    notes = f"""# MaintiQ Predict 최종발표 발표자 노트

## 1. 표지
- MaintiQ Predict는 예지보전 모델 하나가 아니라 CSV 입력, 전처리, 위험 예측, SPC, AI 리포트, 작업지시까지 연결한 제품형 MVP입니다.
- 실제 공장 배포 완료가 아니라, 로컬에서 실행 가능한 데스크톱 시스템과 공개 산업 데이터 검증 결과를 발표합니다.

## 2. 문제 정의
- 사후보전은 고장이 난 뒤 대응하기 때문에 downtime과 정비 비용이 커질 수 있습니다.
- 단일 threshold나 SPC만으로는 복합 센서 조합을 반영하기 어렵습니다.
- 예측 결과가 실제 작업지시와 승인 기록으로 이어지지 않으면 현장 의사결정에 쓰기 어렵습니다.

## 3. 시스템 구조
- 센서 CSV를 넣으면 컬럼 매핑, 품질 진단, 예측, SPC 그래프, AI 리포트, 작업지시 이력까지 이어집니다.
- 운영 이력은 SQLite와 CSV export로 남기고, Admin 콘솔에서 검증 자료를 분리 관리합니다.

## 4. 제품 화면
- Full은 정밀 분석 모드이고, Lite는 빠른 점검 모드입니다.
- Lite는 작은 설치본과 빠른 실행을 위해 경량 점수를 쓰므로 Full 결과와 다를 수 있습니다.

## 5. 데이터 처리
- AI4I 2020 데이터에서 target은 Machine failure입니다.
- UDI, Product ID와 TWF/HDF/PWF/OSF/RNF 누수 컬럼은 제거했습니다.
- Type은 one-hot encoding했고, 회사 CSV는 alias mapping, 단위 변환, 품질 진단을 지원합니다.

## 6. 모델 결과
- XGBoost는 PR-AUC {ev.metrics['models']['xgboost']['pr_auc']:.4f}, ROC-AUC {ev.metrics['models']['xgboost']['roc_auc']:.4f}입니다.
- threshold {th['selected_threshold']:.2f}에서 precision {th['selected_metrics']['precision']:.4f}, recall {th['selected_metrics']['recall']:.4f}, F1 {th['selected_metrics']['f1_score']:.4f}입니다.
- 이 threshold는 F1 기준으로 선택했고, 운영 정책에 따라 바뀔 수 있습니다.

## 7. 경보 정책
- SPC-only는 precision은 높지만 recall이 낮아 놓치는 고장이 많습니다.
- ML threshold는 F1 균형이 가장 좋고, ML+SPC는 recall을 높이는 정책으로 볼 수 있습니다.

## 8. 운영 가치 simulation
- normalized cost simulation은 실제 비용 절감 실증이 아닙니다.
- false alarm, missed failure, planned action에 상대 비용을 부여해 정책 비교 가능성을 평가한 것입니다.

## 9. SCANIA 공개 benchmark
- SCANIA Component X는 실제 fleet 기반 공개 산업 데이터입니다.
- XGBoost official-cost optimized는 rule baseline 대비 official cost metric {best_scania['cost_improvement_vs_rule'] * 100:.2f}% 개선을 보였습니다.
- 단, 이것은 공식 cost matrix 기준 개선이며 실제 회사 원화 비용 절감 실증은 아닙니다.

## 10. 공개 benchmark 확장
- MetroPT-3는 압축기 고장/이상탐지, C-MAPSS는 turbofan RUL, IMS/FEMTO는 bearing run-to-failure 분석에 대응합니다.
- 공개 benchmark는 논문 검증 근거를 보강하지만, 회사별 실제 효과 주장은 별도 로그가 필요합니다.

## 11. 적용 범위
- 구현된 것은 제품형 MVP, 공개 benchmark, field validation 템플릿입니다.
- 실제 PLC/SCADA 운영망, 공장 실시간 배포, 실제 비용 절감률은 아직 주장하지 않습니다.

## 12. 결론
- 새로움은 XGBoost 하나가 아니라, ML+SPC+GenAI+작업지시+검증 근거를 하나의 운영 흐름으로 묶은 점입니다.
- 다음 단계는 실제 회사의 labeled sensor CSV, 정비 이력, downtime/cost 로그 확보입니다.
"""
    SPEAKER_NOTES_PATH.write_text(notes, encoding="utf-8")


def main() -> None:
    ev = load_evidence()
    build_pptx(ev)
    build_previews(ev)
    write_speaker_notes(ev)
    print(f"PPTX saved: {PPTX_PATH}")
    print(f"Preview folder: {PREVIEW_DIR}")
    print(f"Contact sheet: {CONTACT_SHEET_PATH}")
    print(f"Speaker notes: {SPEAKER_NOTES_PATH}")


if __name__ == "__main__":
    main()
